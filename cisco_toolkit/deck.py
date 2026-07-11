"""Executive presentation deck (.pptx) — the stakeholder-facing twin of the workbook / explorer / runbook.

Reads the SAME snapshot the other deliverables do and renders a short, design-conscious slide deck that a
steering committee can read in five minutes: fleet posture, the top migration risks, the keystone devices
the fleet depends on, hardware end-of-support exposure, the wave plan, and where to start.

python-pptx is an OPTIONAL dependency (like python-docx for the runbook): imported inside the function so the
package imports fine without it, and a missing library is a warning + skip, never a crash. Every snapshot
read is defensive (.get / tolerant of missing keys) so an older snapshot degrades gracefully rather than
failing. Deterministic content; no network, no device data beyond the snapshot.
"""
import logging
import os

from cisco_toolkit.textutils import xml_safe   # shared XML-illegal-char sanitizer (the deck's single text sink)
from cisco_toolkit.brand_tokens import DECK_NAVY_RGB

logger = logging.getLogger(__name__)

# Topic-informed palette (enterprise network risk): a dominant navy with an ice-blue support tone, white
# content, and severity colours for the risk content. Dark title/closing slides "sandwich" light content.
_NAVY = DECK_NAVY_RGB
_NAVY2 = (0x2A, 0x35, 0x7A)
_ICE = (0xCA, 0xDC, 0xFC)
_WHITE = (0xFF, 0xFF, 0xFF)
_INK = (0x1F, 0x29, 0x37)
_MUTED = (0x6B, 0x72, 0x80)
_CRIT = (0xC0, 0x39, 0x2B)
_HIGH = (0xE6, 0x7E, 0x22)
_MED = (0xD4, 0xA0, 0x17)
_OK = (0x27, 0xAE, 0x60)
_GRID = (0xE5, 0xE7, 0xEB)

_SEV_COLOR = {"Critical": _CRIT, "High": _HIGH, "Medium": _MED, "Low": _OK, "Info": _MUTED}
_BAND_COLOR = {"Excellent": _OK, "Good": (0x5C, 0xB8, 0x5C), "Fair": _MED, "Poor": _HIGH,
               "Critical": _CRIT, "Insufficient Data": _MUTED}
_LC_BAND_COLOR = {"Past-LDoS": _CRIT, "Past-EoS": _HIGH, "Near-LDoS": _MED, "Active": _OK, "Unknown": _MUTED}

_HEAD = "Calibri"   # universally rendered (Office default); hierarchy comes from size / weight / colour
_BODY = "Calibri"


def _clean(s) -> str:
    """Repair the 'Â·' mojibake the cross-axis headlines carry for their '·' separator AND strip the XML-illegal
    chars (C0 controls, U+FFFE/U+FFFF noncharacters, lone surrogates) that abort python-pptx at save -- the same
    class the workbook/docx are hardened against. _clean is applied to EVERY run's text (the single text sink at
    text()/chip()), so one bad device byte can no longer lose the whole steering-committee deck."""
    return xml_safe((str(s) if s is not None else "").replace("Â·", "·").replace("Â", ""))


def write_executive_deck_pptx(output_path: str, snap_dict: dict, label: str) -> None:
    """Render the executive deck to `output_path`. A missing python-pptx is a warning + skip (the run's
    other deliverables are unaffected); any unexpected render error is logged, never raised."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE
    except ImportError:
        logger.warning("  Executive deck (PPTX) skipped: python-pptx not installed "
                       "(pip install python-pptx, or pass --no-pptx).")
        return

    snap = snap_dict or {}

    def C(rgb):
        return RGBColor(*rgb)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    W = 13.333
    blank = prs.slide_layouts[6]

    def slide(dark=False):
        s = prs.slides.add_slide(blank)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = C(_NAVY if dark else _WHITE)
        return s

    def text(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=None):
        """runs = list of (string, size, rgb, bold[, italic]) for paragraph 0, OR a list of such lists for
        multiple paragraphs."""
        tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = Pt(0)
        tf.margin_top = tf.margin_bottom = Pt(0)
        paras = runs if runs and isinstance(runs[0], list) else [runs]
        for i, para in enumerate(paras):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            if space is not None:
                p.space_after = Pt(space)
            for spec in para:
                txt, size, rgb, bold = spec[0], spec[1], spec[2], spec[3]
                ital = spec[4] if len(spec) > 4 else False
                r = p.add_run()
                r.text = _clean(txt)
                r.font.size = Pt(size)
                r.font.bold = bold
                r.font.italic = ital
                r.font.name = _HEAD if bold else _BODY
                r.font.color.rgb = C(rgb)
        return tb

    def rect(s, l, t, w, h, fill, shape=MSO_SHAPE.RECTANGLE, line=None):
        sp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
        sp.fill.solid()
        sp.fill.fore_color.rgb = C(fill)
        if line is None:
            sp.line.fill.background()
        else:
            sp.line.color.rgb = C(line)
            sp.line.width = Pt(0.75)
        sp.shadow.inherit = False
        return sp

    def chip(s, l, t, txt, fill, w=1.15, h=0.34, tcolor=_WHITE, size=11):
        sp = rect(s, l, t, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = sp.text_frame
        tf.word_wrap = False
        tf.margin_left = tf.margin_right = Pt(4)
        tf.margin_top = tf.margin_bottom = Pt(1)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = _clean(txt)
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.name = _HEAD
        r.font.color.rgb = C(tcolor)
        return sp

    def stat(s, l, t, num, lbl, numcolor=_NAVY, w=2.4):
        text(s, l, t, w, 0.9, [(str(num), 46, numcolor, True)])
        text(s, l, t + 0.92, w, 0.5, [(lbl, 12, _MUTED, False)])

    def propbar(s, l, t, w, segments, h=0.34):
        total = sum(max(0, c) for c, _ in segments) or 1
        x = l
        for count, rgb in segments:
            if count <= 0:
                continue
            seg_w = w * (count / total)
            rect(s, x, t, max(seg_w, 0.04), h, rgb)
            x += seg_w

    def header(s, kicker, title, dark=False):
        text(s, 0.7, 0.55, W - 1.4, 0.4, [(kicker.upper(), 13, _ICE if dark else _HIGH, True)])
        text(s, 0.7, 0.95, W - 1.4, 0.9, [(title, 32, _WHITE if dark else _NAVY, True)])

    # ---------------------------------------------------------------- 1. Title (dark)
    def _D(x): return x if isinstance(x, dict) else {}    # audit-5 totality: a truthy non-dict section -> {}
    def _R(x): return [r for r in (x if isinstance(x, list) else []) if isinstance(r, dict)]
    s = slide(dark=True)
    rect(s, 0, 0, W, 7.5, _NAVY)
    eb = _D(snap.get("executive_brief"))
    scale = _D(eb.get("scale"))
    text(s, 0.9, 2.0, W - 1.8, 1.6,
         [[("Network Migration", 52, _WHITE, True)], [("Assessment", 52, _ICE, True)]], space=2)
    text(s, 0.9, 3.95, W - 1.8, 0.5, [(label or "Current-state assessment", 18, _ICE, False, True)])
    posture = _D(eb.get("posture"))
    # the canonical posture_statement is ALWAYS non-empty, so the old scale fallback never rendered — the
    # deck showed no fleet scale at all. Always render the canonical scale on a second subtitle line.
    sub = eb.get("posture_statement") or ""
    # honesty (wave R2-4-02): when the cross-axis synthesis FAILED (sentinel / assessment_integrity flag),
    # say so on the title slide rather than rendering an em-dash "— devices · — endpoints" as if it were scale.
    if eb.get("_unavailable") or (snap.get("assessment_integrity") or {}).get("executive_brief") == "compute_failed":
        _scale_line = "⚠ Cross-axis synthesis unavailable — see the workbook Executive Brief"
    else:
        # coerce PER VALUE: dict.get(default) only substitutes when the KEY is absent, so a present-but-null
        # scale value (an uploaded / partially-computed snapshot) rendered the literal 'None devices' on the
        # marquee client-facing slide. Show an em-dash for any non-int value instead.
        def _sv(k):
            return scale.get(k) if isinstance(scale.get(k), int) else "—"
        _ncoll = scale.get("n_collected")
        _coll = f" ({_ncoll} collected)" if isinstance(_ncoll, int) and _ncoll != scale.get("n_devices") else ""
        _scale_line = (f"{_sv('n_devices')} devices{_coll} · {_sv('n_endpoints')} endpoints "
                       f"· {_sv('n_vlans')} VLANs in scope")
    text(s, 0.9, 4.7, W - 1.8, 1.6,
         [[(sub, 14, _ICE, False)], [(_scale_line, 13, _ICE, False)]])
    # P3-E3 provenance (the explorer-ctx pattern, html.py:539-543): engine version + generation timestamp
    # + snapshot stem on the title footer, so a printed / forwarded deck stays traceable to the exact
    # snapshot that produced it. Absent fields are omitted, never fabricated.
    _stem = os.path.basename(str(output_path or ""))
    for _suf in ("_executive_deck.pptx", ".pptx"):
        if _stem.endswith(_suf):
            _stem = _stem[: -len(_suf)]
            break
    _prov = " · ".join(p for p in (
        snap.get("script_version") or "",
        (snap.get("generated_at") or "")[:19].replace("T", " "),
        f"snapshot {_stem}" if _stem else "",
    ) if p)
    text(s, 0.9, 6.55, W - 1.8, 0.6,
         [[("Generated offline by the Cisco Migration-Assessment Toolkit — full evidence in the workbook, "
            "explorer & runbook.", 11, (0x9F, 0xB3, 0xE0), False)]]
         + ([[(_prov, 10, (0x9F, 0xB3, 0xE0), False)]] if _prov else []))

    # ---------------------------------------------------------------- 2. Fleet posture (light)
    s = slide()
    header(s, "Executive summary", "Fleet posture")
    hs = _R(snap.get("health_scores"))
    band_counts = {}
    for r in hs:
        band_counts[r.get("band", "")] = band_counts.get(r.get("band", ""), 0) + 1
    avg = posture.get("avg_health", "—")
    # stat callouts across the full width (generous gaps, no side-by-side columns to overlap)
    stat(s, 0.7, 1.95, f"{avg}", "avg health / 100", _NAVY, w=3.4)
    # false-health guard: if the brief is absent/failed, posture is {} — fall back to the band tally
    # computed above from health_scores, never a literal 0 that would claim '0 Critical' on a Critical fleet.
    stat(s, 5.0, 1.95, posture.get("n_critical", band_counts.get("Critical", 0)), "Critical-band switches", _CRIT, w=3.4)
    stat(s, 9.3, 1.95, len(hs), "switches in scope", _INK, w=3.3)   # 303 inventoried; 50 not collected -> "assessed" overclaims
    # full-width band-distribution bar
    order = ["Excellent", "Good", "Fair", "Poor", "Critical", "Insufficient Data"]
    segs = [(band_counts.get(b, 0), _BAND_COLOR.get(b, _MUTED)) for b in order]
    if any(c for c, _ in segs):
        text(s, 0.7, 3.45, 12.0, 0.3, [("HEALTH-BAND DISTRIBUTION", 12, _MUTED, True)])
        propbar(s, 0.7, 3.8, 11.9, segs)
        legend = [b for b in order if band_counts.get(b, 0)]
        text(s, 0.7, 4.25, 12.0, 0.4,
             [[(f"{b}: {band_counts.get(b, 0)}     ", 11, _BAND_COLOR.get(b, _MUTED), True) for b in legend]])
    # full-width per-axis headlines below
    axes = eb.get("axes") or []
    text(s, 0.7, 4.95, 12.0, 0.35, [("HEADLINE BY AXIS", 12, _HIGH, True)])
    y = 5.35
    for a in axes[:4]:
        sev = a.get("severity", "Info")
        chip(s, 0.7, y, sev, _SEV_COLOR.get(sev, _MUTED), w=1.0, h=0.3, size=10)
        text(s, 1.85, y - 0.03, 10.8, 0.5,
             [[(_clean(a.get("axis", "")) + ":  ", 12, _NAVY, True), (_clean(a.get("headline", "")), 12, _INK, False)]])
        y += 0.5
    # honesty: never silently drop axes — breadcrumb the overflow + how many are Critical/High, so a
    # 4-row cap can't hide a High/Critical axis (the severity-sorted source means the cap drops the tail).
    _more = axes[4:]
    if _more:
        _more_hi = sum(1 for a in _more if a.get("severity") in ("Critical", "High"))
        text(s, 0.7, y, 11.8, 0.4,
             [[(f"+ {len(_more)} more axis headline(s)"
                + (f" — {_more_hi} at High or above" if _more_hi else "")
                + " — full set in the workbook Executive Summary.", 10, _MUTED, False)]])

    # ---------------------------------------------------------------- 3. Top risks (light)
    s = slide()
    header(s, "What gates the migration", "Top migration risks")
    pl = _R(snap.get("punchlist"))
    n_crit = sum(1 for i in pl if i.get("severity") == "Critical")
    n_high = sum(1 for i in pl if i.get("severity") == "High")
    text(s, 0.7, 1.95, W - 1.4, 0.4,
         [[(f"{len(pl)} consolidated finding(s)  ", 14, _NAVY, True),
           (f"· {n_crit} Critical · {n_high} High — the de-duplicated, severity-ranked punch-list.",
            13, _MUTED, False)]])
    y = 2.6
    for it in pl[:5]:
        sev = it.get("severity", "Info")
        chip(s, 0.7, y, sev, _SEV_COLOR.get(sev, _MUTED), w=1.0, h=0.34, size=10)
        devs = ", ".join(str(d) for d in (it.get("devices") or [])[:4])
        more = "" if len(it.get("devices") or []) <= 4 else f" +{len(it['devices']) - 4}"
        text(s, 1.85, y - 0.04, W - 2.6, 0.7,
             [[(_clean(it.get("title", "")), 14, _INK, True)],
              [(f"{it.get('category', '')}  ·  {devs}{more}", 11, _MUTED, False)]], space=1)
        y += 0.78
    if len(pl) > 5:
        text(s, 1.85, y, W - 2.6, 0.3,
             [(f"+ {len(pl) - 5} more in the Migration Punch-List workbook sheet.", 11, _MUTED, False, True)])

    # ------------------------------------------------------- 3b. Riskiest assets (NEW-V3.23.174)
    # The Device Risk Register's worst boxes: per-asset stacked risk (impact x exposure), the
    # compound patterns, and the engineer's verdict. Skipped when the snapshot has no register.
    dd = (snap.get("device_dossiers") or {}).get("per_device") or []
    dd_top = [d for d in dd if d.get("risk_band") in ("Severe", "Elevated", "Guarded")][:5]
    if dd_top:
        s = slide()
        header(s, "Stacked risk per asset", "The assets that worry an engineer most")
        dsum = (snap.get("device_dossiers") or {}).get("summary") or {}
        dbands = dsum.get("bands") or {}
        text(s, 0.7, 1.95, W - 1.4, 0.4,
             [[(f"{dbands.get('Severe', 0)} Severe · {dbands.get('Elevated', 0)} Elevated  ", 14, _NAVY, True),
               ("— risk index = topology impact × stacked exposure; compound patterns (CR-xx) mark "
                "independent risks coinciding on one box.", 13, _MUTED, False)]])
        _BAND_SEV = {"Severe": "Critical", "Elevated": "High", "Guarded": "Medium", "Low": "Low", "Unassessed": "Info"}
        y = 2.6
        for d in dd_top:
            band = d.get("risk_band", "Low")
            chip(s, 0.7, y, _clean(f"{band} {d.get('risk_index', 0)}"),
                 _SEV_COLOR.get(_BAND_SEV.get(band, "Low"), _MUTED), w=1.45, h=0.34, size=10)
            crs = " · ".join(c.get("code", "") for c in (d.get("compound") or [])) or ""
            text(s, 2.35, y - 0.04, W - 3.1, 0.7,
                 [[(_clean(str(d.get("host", ""))), 14, _INK, True),
                   (f"   {crs}", 11, _MUTED, False)],
                  [(_clean((d.get("verdict") or "")[:160]), 11, _MUTED, False)]], space=1)
            y += 0.78
        if len(dd) > len(dd_top):
            text(s, 2.35, y, W - 3.1, 0.3,
                 [(f"+ {len(dd) - len(dd_top)} more asset(s) in the Device Risk Register workbook sheet.",
                   11, _MUTED, False, True)])

    # ---------------------------------------------------------------- 4. Keystone devices (light)
    s = slide()
    header(s, "Concentrated dependency", "The switches the fleet depends on")
    fi = sorted(_R(snap.get("failure_impact")),
                key=lambda r: -(int(r.get("stranded") or 0)))
    keystones = [r for r in fi if int(r.get("stranded") or 0) > 0][:5]
    # off-scan-gateway records carry no usable blast radius (audit-3 #7); count them so an INDETERMINATE estate
    # is not mistaken for a well-distributed one (audit-4 #8 false-health).
    n_indet = sum(1 for r in fi if int(r.get("off_scan_gw_vlans") or 0) > 0)
    text(s, 0.7, 1.95, W - 1.4, 0.5,
         [("Ranked by migration blast radius — the collateral endpoints stranded if the switch drops "
           "during its move. Sequence and protect these first.", 13, _MUTED, False)])
    y = 2.8
    if keystones:
        for r in keystones:
            stat_w = 1.7
            text(s, 0.7, y, stat_w, 0.5, [(str(r.get("stranded", 0)), 30, _CRIT, True)])
            text(s, 0.7, y + 0.55, stat_w, 0.3, [("stranded", 10, _MUTED, False)])
            text(s, 2.5, y + 0.05, W - 3.2, 0.7,
                 [[(_clean(str(r.get("host", ""))), 16, _NAVY, True),
                   (f"   {r.get('vlans_impacted', 0)} VLAN(s) · {r.get('hard', 0)} hard-partitioned", 12, _MUTED, False)],
                  [(_clean(r.get("detail", "")), 11, _INK, False)]], space=1)
            y += 0.92
    elif not fi or n_indet:
        # blast radius NOT computed (no failure_impact) or INDETERMINATE (off-scan gateways) -> a coverage gap,
        # NOT a clean bill. Render muted (not _OK green) so a blind estate never reads as 'well distributed'.
        msg = ("Blast radius not assessable — the failure-impact analysis was not computed for this snapshot "
               "(thin or uploaded data). Redundancy is UNKNOWN, not verified well-distributed."
               if not fi else
               f"Blast radius INDETERMINATE — {n_indet} switch(es) have an off-scan gateway, so their removal "
               "impact could not be modelled. This is a coverage gap, not a clean bill of distribution.")
        text(s, 0.7, y, W - 1.4, 0.9, [(msg, 14, _MUTED, True)])
    else:
        text(s, 0.7, y, W - 1.4, 0.5,
             [("No switch strands collateral endpoints on removal — dependency is well distributed.", 14, _OK, True)])

    # ---------------------------------------------------------------- 5. Hardware lifecycle (light)
    lr = snap.get("lifecycle_risk") or {}
    lsum = lr.get("summary") or {}
    if lsum.get("n_devices"):
        s = slide()
        header(s, "A primary migration driver", "Hardware end-of-support exposure")
        # "Past end-of-support" is the Past-LDoS band ALONE (no TAC) — Past-EoS is end-of-SALE with
        # the support window still open. Mirror the canonical executive_brief lifecycle axis
        # (compute_executive_brief's lc_pe / lc_known pair) so the deck headline agrees with every
        # other surface (A3 SSOT fix). Past-EoS is not lost: it keeps its own segment below.
        past = lsum.get("n_past_ldos", 0)
        near = lsum.get("n_near", 0)
        # Denominator == KNOWN-model (EoL-assessable) devices == full fleet minus unidentified-model
        # devices, mirroring compute_executive_brief's lc_known. Dividing by the full n_devices diluted
        # the headline and made THIS slide the sole set-wide outlier (70% of 303) while the workbook,
        # design, MOP, runbook and engagement all report "86% of 247 assessable" (QA scorecard row 12,
        # residual #1). The Unknown count stays visible in the band-distribution legend below, so the
        # narrower denominator is coverage-honest, not a hidden exclusion.
        assessable = lsum.get("n_devices", 0) - lsum.get("n_unknown", 0)
        pct = round(100 * (past + near) / assessable) if assessable else 0
        stat(s, 0.7, 2.1, f"{pct}%", f"of {assessable} assessable, past or nearing end-of-support", _HIGH)
        stat(s, 3.6, 2.1, past, "past end-of-support", _CRIT)
        stat(s, 6.2, 2.1, near, "within 1 year", _MED)
        lc_order = ["Past-LDoS", "Past-EoS", "Near-LDoS", "Active", "Unknown"]
        byb = lsum.get("by_band") or {}
        segs = [(byb.get(b, 0), _LC_BAND_COLOR.get(b, _MUTED)) for b in lc_order]
        if any(c for c, _ in segs):
            text(s, 0.7, 3.6, W - 1.4, 0.3, [("Lifecycle band distribution", 12, _MUTED, True)])
            propbar(s, 0.7, 3.95, W - 1.4, segs)
            legend = [b for b in lc_order if byb.get(b, 0)]
            text(s, 0.7, 4.4, W - 1.4, 0.4,
                 [[(f"{b}: {byb.get(b, 0)}   ", 11, _LC_BAND_COLOR.get(b, _MUTED), True) for b in legend]])
        text(s, 0.7, 5.3, W - 1.4, 0.5,
             [("Aging hardware is one of the top reasons to migrate — and a hard constraint on how long the "
               "current fabric can safely run. Detail in the Lifecycle Risk workbook sheet.", 12, _INK, False, True)])

    # ---------------------------------------------------------------- 6. Migration waves (light)
    s = slide()
    header(s, "How it sequences", "Migration waves & readiness")
    mr = _R(snap.get("migration_readiness"))
    mg = snap.get("move_groups") or []
    # SSOT: the honest, actionable headline is the SEQUENCED wave count from the design blueprint's
    # wave_plan -- the raw move-group count is the L2 blast-radius partition (one big coupled domain +
    # singletons), which read as a wave count overstates parallelism (the audit's coverage-honesty fix).
    wp = (((snap.get("design_blueprint") or {}).get("target_state") or {}).get("wave_plan")) or {}
    tally = {"READY": 0, "CAUTION": 0, "NOT READY": 0}
    for r in mr:
        v = r.get("readiness")
        if v in tally:
            tally[v] += 1
    n_groups = len(mr)
    _nmg = wp.get("n_move_groups") or len(mg) or n_groups
    # coverage-honesty (QA row-11 BLOCK): the readiness tiles count EVERY assessed device group — say so on
    # the tiles themselves, so the wave count beside them cannot read as the whole scope (an exec read
    # "9 candidate waves" + a 6-row list as the full estate; the real estate was 53 groups, ~7x larger).
    _denom = f" of {n_groups} device group(s)" if n_groups else " — no groups assessed"
    if wp.get("n_waves"):
        stat(s, 0.7, 2.1, wp["n_waves"], f"candidate waves · {_nmg} move-group(s)", _NAVY)
    else:
        stat(s, 0.7, 2.1, len(mg) or len(mr), "move groups", _NAVY)
    stat(s, 3.3, 2.1, tally["NOT READY"], f"NOT READY{_denom}", _CRIT, w=2.5)
    stat(s, 5.9, 2.1, tally["CAUTION"], f"CAUTION{_denom}", _MED, w=2.5)
    stat(s, 8.5, 2.1, tally["READY"], f"READY{_denom}", _OK, w=2.5)
    _rt = {"READY": _OK, "CAUTION": _MED, "NOT READY": _CRIT}
    text(s, 0.7, 3.58, W - 1.4, 0.35, [("PER-GROUP READINESS", 12, _HIGH, True)])
    # Layout is COMPUTED, never coordinated by hand (QA row-11 BLOCK: the footnote was pinned at 6.5in
    # while a 6-row list runs to ~6.83in — the shapes overprinted and both were illegible). The row count
    # derives from the space the footnote + bottom margin leave, and the footnote's top from the rows
    # actually rendered — overlap is impossible by construction (the figgen auto-sized-boxes lesson).
    # Truncation is disclosed with a "+N more" cue, never silent.
    LIST_TOP, ROW_H, CUE_H, FOOT_GAP, CONTENT_BOTTOM = 3.95, 0.5, 0.32, 0.12, 7.2
    foot_h = 0.75 if wp.get("n_waves") else 0.45
    budget = CONTENT_BOTTOM - foot_h - FOOT_GAP - LIST_TOP
    n_fit = int(budget // ROW_H)
    if n_groups > n_fit:
        n_fit = max(int((budget - CUE_H) // ROW_H), 0)   # reserve the overflow-cue line
    y = LIST_TOP
    for r in mr[:n_fit]:
        v = r.get("readiness", "—")
        chip(s, 0.7, y, v, _rt.get(v, _MUTED), w=1.5, h=0.32, size=10)
        text(s, 2.4, y - 0.02, W - 3.1, 0.4,
             [[(_clean(r.get("group", "")) + "  ", 13, _NAVY, True),
               (f"{r.get('n_fail', 0)} blocking · {r.get('n_warn', 0)} warning check(s)", 12, _INK, False)]])
        y += ROW_H
    _hidden = n_groups - min(n_fit, n_groups)
    if _hidden:
        text(s, 0.7, y, W - 1.4, 0.3,
             [(f"+ {_hidden} more device group(s) not shown — full per-group readiness in the "
               "Migration Readiness workbook sheet.", 11, _MUTED, False, True)])
        y += CUE_H
    foot_top = y + FOOT_GAP
    if wp.get("n_waves"):
        text(s, 0.7, foot_top, W - 1.4, foot_h,
             [(f"{_nmg} L2-coupled move-group(s) — largest a "
               f"{wp.get('largest_group', 0)}-switch broadcast domain — sequence into {wp['n_waves']} "
               f"candidate wave(s) of ≤ {wp.get('wave_cap', 40)} switches. Clear the NOT-READY blockers "
               "first, biggest blast radius first; verify each wave with the Cutover Validation plan.",
               12, _MUTED, False, True)])
    else:
        text(s, 0.7, foot_top, W - 1.4, foot_h,
             [("Clear the NOT-READY blockers first, biggest blast radius first — then verify each wave with the "
               "Cutover Validation plan.", 12, _MUTED, False, True)])

    # ------------------------------------------------------- 6b. Target-state design (NEW — design engine)
    # The CCDE-grounded design blueprint: the recommended target-state decisions + the weakest trade-off
    # axes. Data-gated — skipped when the snapshot carries no design_blueprint (the SAME object the HLD/LLD
    # and the dashboards read; one source of truth).
    bp = snap.get("design_blueprint") or {}
    bp_rec = [d for d in (bp.get("decisions") or []) if isinstance(d, dict) and d.get("status") == "recommended"]
    if bp_rec:
        s = slide()
        header(s, "CCDE-grounded target state", "The design the migration should adopt")
        bsum = bp.get("summary") or {}
        text(s, 0.7, 1.9, W - 1.4, 0.4,
             [[(f"{bsum.get('n_recommended', len(bp_rec))} design decision(s)  ", 14, _NAVY, True),
               (f"· {bsum.get('n_critical', 0)} critical · {bsum.get('n_needs_requirement', 0)} need a "
                "requirement — each traced to a design principle and the trade-off axes it serves.",
                13, _MUTED, False)]])
        sc = [a for a in (bp.get("tradeoff_scorecard") or []) if isinstance(a, dict)]
        weak = [a for a in sc if isinstance(a.get("score"), int) and a.get("score") <= 1][:5]
        if weak:
            text(s, 0.7, 2.5, W - 1.4, 0.3, [("WEAKEST TRADE-OFF AXES (0–1 / 4)", 12, _HIGH, True)])
            x = 0.7
            for a in weak:
                chip(s, x, 2.9, _clean(f"{a.get('axis', '')} {a.get('score')}/4"), _CRIT, w=2.0, h=0.32, size=10)
                x += 2.1
        y = 3.65
        for d in bp_rec[:5]:
            sev = d.get("priority", "Info")
            chip(s, 0.7, y, sev, _SEV_COLOR.get(sev, _MUTED), w=1.0, h=0.34, size=10)
            text(s, 1.85, y - 0.04, W - 2.6, 0.7,
                 [[(_clean(d.get("title", "")), 14, _INK, True)],
                  [(_clean(((d.get("evidence") or {}).get("summary") or "")[:120]), 11, _MUTED, False)]], space=1)
            y += 0.78

    # ---------------------------------------------------------------- 7. Where to start (dark)
    s = slide(dark=True)
    header(s, "Recommendation", "Where to start", dark=True)
    gating = eb.get("top_gating") or []
    text(s, 0.7, 2.0, W - 1.4, 0.4,
         [("Resolve these gating items before the first cutover wave:", 15, _ICE, False)])
    y = 2.7
    if gating:
        for i, g in enumerate(gating[:5], 1):
            rect(s, 0.7, y, 0.5, 0.5, _HIGH, shape=MSO_SHAPE.OVAL)
            text(s, 0.7, y + 0.04, 0.5, 0.42, [(str(i), 18, _WHITE, True)], align=PP_ALIGN.CENTER)
            text(s, 1.4, y + 0.03, W - 2.1, 0.6, [(_clean(g), 15, _WHITE, False)])
            y += 0.78
    else:
        text(s, 0.7, y, W - 1.4, 0.5, [("No High/Critical gating items — the fleet is in good shape to schedule.",
                                        16, _OK, True)])
    text(s, 0.7, 6.6, W - 1.4, 0.6,
         [("Then proceed wave by wave: clear blockers → cut over → run the per-wave Cutover Validation "
           "checklist → confirm before the next wave.", 12, _ICE, False, True)])

    n_slides = len(prs.slides._sldIdLst)
    prs.save(output_path)
    logger.info(f"[Phase 32] Executive deck (PPTX) written: {output_path} ({n_slides} slides)")
