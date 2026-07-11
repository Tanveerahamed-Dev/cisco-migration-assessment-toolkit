"""NEW-V3.23.144: executive presentation deck (.pptx). Pins the slide count, the signature content drawn
from each snapshot axis, the 'Â·' mojibake cleanup, and graceful degradation on a sparse snapshot.
python-pptx is an optional dependency — these tests skip when it isn't installed (CI installs it)."""
import pytest

pytest.importorskip("pptx", reason="python-pptx not installed (optional deck dependency)")

from pptx import Presentation                                   # noqa: E402
from cisco_toolkit.deck import write_executive_deck_pptx        # noqa: E402


def _rich_snap():
    return {
        "executive_brief": {
            "scale": {"n_devices": 12, "n_domains": 3, "n_endpoints": 840},
            "posture": {"avg_health": 61, "n_critical": 2, "n_poor": 1, "worst_band": "Critical"},
            "posture_statement": "Migration posture: 2 switches Critical; hardware EoL a primary driver.",
            "axes": [
                {"axis": "Fleet health", "severity": "Critical", "headline": "61/100 avg · 2 Critical"},
                {"axis": "Punch-list", "severity": "High", "headline": "30 items · 2 Critical, 11 High"},
            ],
            "top_gating": ["61/100 avg · 2 Critical", "30 items · 2 Critical, 11 High", "40% past/near EoS"],
        },
        "health_scores": [
            {"switch": "core1", "score": 12, "band": "Critical"},
            {"switch": "core2", "score": 55, "band": "Fair"},
            {"switch": "acc1", "score": 88, "band": "Excellent"},
        ],
        "punchlist": [
            {"severity": "Critical", "category": "Cross-layer", "title": "VLAN 30 single-fiber uplink",
             "devices": ["core1", "acc1"]},
            {"severity": "High", "category": "STP", "title": "Accidental root on VLAN 10", "devices": ["acc1"]},
        ] + [{"severity": "Low", "category": "Hygiene", "title": f"unused acl {i}", "devices": ["x"]}
             for i in range(8)],
        "failure_impact": [
            {"host": "core1", "severity": "High", "vlans_impacted": 4, "stranded": 220, "hard": 180,
             "detail": "VLAN 10 hard partition"},
            {"host": "core2", "severity": "Medium", "vlans_impacted": 2, "stranded": 40, "hard": 0,
             "detail": "backup-covered"},
        ],
        "lifecycle_risk": {"summary": {"n_devices": 12, "by_band": {"Past-EoS": 3, "Near-LDoS": 2, "Active": 7},
                                       "n_past_eos": 3, "n_past_ldos": 0, "n_near": 2, "n_active": 7,
                                       "n_unknown": 0}},
        "migration_readiness": [
            {"group": "Group 1", "readiness": "NOT READY", "n_fail": 3, "n_warn": 2},
            {"group": "Group 2", "readiness": "READY", "n_fail": 0, "n_warn": 0},
        ],
        "move_groups": [{"switches": ["core1", "acc1"]}, {"switches": ["core2"]}],
    }


def _deck(path):
    p = Presentation(path)
    n = len(p.slides._sldIdLst)
    txt = "\n".join(sh.text_frame.text for sl in p.slides for sh in sl.shapes if sh.has_text_frame)
    return n, txt


def test_deck_title_slide_always_shows_canonical_scale(tmp_path):
    """The title slide must ALWAYS surface the canonical fleet scale (devices/endpoints/VLANs). Previously
    scale rendered only as a FALLBACK for an empty posture_statement — which is canonically never empty —
    so the deck showed no fleet scale at all."""
    out = tmp_path / "deck_scale.pptx"
    snap = _rich_snap()
    snap["executive_brief"]["scale"] = {"n_devices": 303, "n_endpoints": 5127, "n_vlans": 202}
    write_executive_deck_pptx(str(out), snap, "Test fleet")
    _, txt = _deck(str(out))
    assert "303 devices" in txt and "5127 endpoints" in txt and "202 VLANs" in txt


def test_deck_discloses_failed_brief_not_false_zeros(tmp_path):
    """R2-4-02: a FAILED cross-axis brief (assessment_integrity flag / _unavailable sentinel) must be
    DISCLOSED on the deck, and the Critical-band count must still come from health_scores -- never a silent
    em-dash scale line nor a false '0 Critical'."""
    snap = {
        "executive_brief": {"_unavailable": True},
        "assessment_integrity": {"executive_brief": "compute_failed"},
        "health_scores": [
            {"switch": "core1", "score": 8, "band": "Critical"},
            {"switch": "core2", "score": 11, "band": "Critical"},
            {"switch": "acc1", "score": 90, "band": "Excellent"},
        ],
        "punchlist": [], "failure_impact": [], "migration_readiness": [], "move_groups": [],
        "lifecycle_risk": {"summary": {}},
    }
    out = tmp_path / "deck_failed_brief.pptx"
    write_executive_deck_pptx(str(out), snap, "Fleet X")
    _n, txt = _deck(str(out))
    assert "unavailable" in txt.lower()      # the synthesis failure is disclosed on the title slide
    assert "Critical: 2" in txt              # the REAL Critical count surfaces from the band tally, not 0
    assert "— devices" not in txt            # no silent em-dash scale line masquerading as fleet scale


def test_deck_has_all_slides_and_key_content(tmp_path):
    out = tmp_path / "deck.pptx"
    write_executive_deck_pptx(str(out), _rich_snap(), "Test fleet")
    assert out.is_file()
    n, txt = _deck(str(out))
    assert n == 7, f"expected 7 slides, got {n}"
    assert "Network Migration" in txt                                  # 1 title
    assert "Fleet posture" in txt and "61" in txt                      # 2 posture + avg health
    assert "Top migration risks" in txt and "VLAN 30 single-fiber uplink" in txt   # 3 risks
    assert "core1" in txt and "220" in txt                             # 4 keystone + stranded count
    assert "end-of-support" in txt.lower()                             # 5 lifecycle
    assert "NOT READY" in txt and "Group 1" in txt                     # 6 waves
    assert "Where to start" in txt                                     # 7 recommendation


def test_deck_lifecycle_past_end_of_support_is_ldos_not_eos(tmp_path):
    """A3 (SSOT/coverage-honesty): the 'past end-of-support' headline must read n_past_ldos ALONE
    (matching the canonical executive_brief lifecycle axis, analyze.py:5018), NOT n_past_eos +
    n_past_ldos. Past-EoS is end-of-SALE (support window still open). Fixture: 152 LDoS + 40 EoS +
    61 near of 303 → headline 152 and 70% past/nearing; the old conflation rendered 192 / 83%."""
    snap = _rich_snap()
    snap["lifecycle_risk"] = {"summary": {
        "n_devices": 303, "n_past_ldos": 152, "n_past_eos": 40, "n_near": 61, "n_active": 50,
        "by_band": {"Past-LDoS": 152, "Past-EoS": 40, "Near-LDoS": 61, "Active": 50}}}
    out = tmp_path / "deck.pptx"
    write_executive_deck_pptx(str(out), snap, "Test fleet")
    n, txt = _deck(str(out))
    assert "152" in txt and "70%" in txt           # n_past_ldos headline + (152+61)/303 pct
    assert "192" not in txt and "83%" not in txt    # old n_past_eos+n_past_ldos conflation gone


def test_deck_gains_riskiest_assets_slide_with_register(tmp_path):
    """NEW-V3.23.174: a snapshot carrying the Device Risk Register renders the extra
    'riskiest assets' slide (8 total); the back-compat 7-slide pin above proves the
    slide is data-gated, never empty filler."""
    snap = _rich_snap()
    snap["device_dossiers"] = {
        "per_device": [
            {"host": "core1", "risk_band": "Severe", "risk_index": 60, "impact_score": 10,
             "exposure_score": 6,
             "compound": [{"code": "CR-01", "title": "End-of-support keystone", "severity": "Critical",
                           "basis": "EoL x blast radius"}],
             "verdict": "Stabilize or replace before migration — open advisory surface."},
            {"host": "acc1", "risk_band": "Low", "risk_index": 2, "impact_score": 2,
             "exposure_score": 1, "compound": [], "verdict": "No stacked risk."},
        ],
        "summary": {"n_devices": 2, "bands": {"Severe": 1, "Elevated": 0, "Guarded": 0, "Low": 1},
                    "n_compound": 1, "worst": ["core1"]}}
    out = tmp_path / "deck_rr.pptx"
    write_executive_deck_pptx(str(out), snap, "Test fleet")
    n, txt = _deck(str(out))
    assert n == 8, f"expected 8 slides with the register, got {n}"
    assert "worry an engineer most" in txt
    assert "CR-01" in txt and "Stabilize or replace" in txt


def test_deck_gains_target_state_design_slide(tmp_path):
    """NEW: a snapshot carrying the design_blueprint renders the extra 'target-state design' slide (8
    total); the back-compat 7-slide pin above proves the slide is data-gated (the design engine's
    compute_design_blueprint), never empty filler — the SAME blueprint behind the HLD/LLD and dashboards."""
    snap = _rich_snap()
    snap["design_blueprint"] = {
        "summary": {"n_decisions": 2, "n_recommended": 2, "n_needs_requirement": 1, "n_critical": 1,
                    "headline": "2 design decisions."},
        "tradeoff_scorecard": [{"axis": "availability", "label": "High availability", "score": 0,
                                "posture": "Weak", "evidence": "no FHRP"}],
        "decisions": [
            {"id": "fhrp-first-hop-gateway-redundancy", "title": "Introduce first-hop redundancy",
             "priority": "Critical", "status": "recommended", "evidence": {"summary": "52 VLANs without FHRP"},
             "principle": {"citation": "CCDE In Depth — HA"}, "recommended_action": "HSRP/VRRP"},
            {"id": "x", "title": "Right-size availability", "priority": "High", "status": "needs-requirement",
             "evidence": {"summary": ""}, "principle": {"citation": "CCDE"},
             "requirements_needed": ["availability_tier"]},
        ],
        "coverage": {"caveat": "grounded only in collected evidence"},
    }
    out = tmp_path / "deck_design.pptx"
    write_executive_deck_pptx(str(out), snap, "Test fleet")
    n, txt = _deck(str(out))
    assert n == 8, f"expected 8 slides with the design blueprint, got {n}"
    assert "design the migration should adopt" in txt.lower() or "target state" in txt.lower()
    assert "Introduce first-hop redundancy" in txt


def test_deck_migration_slide_surfaces_honest_wave_count(tmp_path):
    """B1 (audit fix): when the snapshot carries the design wave_plan, the 'How it sequences' slide must
    headline the honest SEQUENCED wave count (design_blueprint.target_state.wave_plan.n_waves) -- not the
    raw move-group count presented as if it were parallelizable waves. A 60-switch L2-coupled domain is
    ONE set sliced into sequenced sub-waves, not 6 parallel waves; the slide must say so."""
    snap = _rich_snap()
    snap["move_groups"] = ([{"switches": [f"s{i}" for i in range(60)]}]
                           + [{"switches": [f"t{j}"]} for j in range(5)])   # 1x60 + 5 singletons = 6 groups
    snap["design_blueprint"] = {
        "summary": {"n_decisions": 1, "n_recommended": 1, "n_needs_requirement": 0, "n_critical": 1,
                    "headline": "1 critical recommended."},
        "tradeoff_scorecard": [], "coverage": {},
        "decisions": [{"id": "fhrp-first-hop-gateway-redundancy", "title": "Introduce FHRP",
                       "priority": "Critical", "status": "recommended",
                       "evidence": {"summary": "x"}, "principle": {"citation": "CCDE"}}],
        "target_state": {"wave_plan": {"n_waves": 3, "n_move_groups": 6, "largest_group": 60,
                                       "wave_cap": 40, "waves": [], "note": "n"}},
    }
    out = tmp_path / "deck_waves.pptx"
    write_executive_deck_pptx(str(out), snap, "Test fleet")
    _n, txt = _deck(str(out))
    low = txt.lower()
    assert "candidate wave" in low, "slide must surface the honest sequenced wave count label"
    # the honest relationship (move-groups -> sequenced waves) must be disclosed, not just the raw 6
    assert "sequence" in low and "60" in txt, "must disclose the largest L2 domain sequences into waves"


def test_deck_cleans_mojibake(tmp_path):
    snap = {"executive_brief": {"axes": [{"axis": "X", "severity": "High", "headline": "a Â· b"}],
                                "top_gating": []}}
    out = tmp_path / "d.pptx"
    write_executive_deck_pptx(str(out), snap, "L")
    _n, txt = _deck(str(out))
    assert "Â·" not in txt and "a · b" in txt


def test_sparse_snapshot_is_tolerated(tmp_path):
    out = tmp_path / "empty.pptx"
    write_executive_deck_pptx(str(out), {}, "Empty")               # no keys at all
    assert out.is_file()
    n, _ = _deck(str(out))
    assert n == 6        # every slide renders except the data-gated lifecycle slide


def _deck_text(path):
    prs = Presentation(path)
    return "\n".join(sh.text_frame.text for sl in prs.slides for sh in sl.shapes if sh.has_text_frame)


def test_deck_survives_xml_illegal_chars(tmp_path):
    """[audit-4 #2 totality] a U+FFFE/U+FFFF/lone-surrogate in any device-derived string aborted the whole
    executive deck at python-pptx save -- the class the workbook/docx are hardened against. _clean (the single
    text sink) must strip them."""
    snap = _rich_snap()
    bad = chr(0xFFFF) + chr(0xFFFE) + chr(0xD800)
    snap["failure_impact"][0]["host"] = snap["failure_impact"][0]["host"] + bad
    snap["executive_brief"]["axes"][0]["headline"] = "61/100 " + bad
    out = str(tmp_path / "d.pptx")
    write_executive_deck_pptx(out, snap, "AJ" + bad)     # must not raise
    Presentation(out)                                    # and open


def _m0_like_snap():
    """Mirrors the real M0 estate that exposed the QA row-11 BLOCK (scorecard 2026-07-11): 53 device
    groups (1 NOT READY / 52 CAUTION / 0 READY) and a wave_plan of 9 candidate waves over 53 move-groups,
    largest a 251-switch domain — the shape on which the fixed 6.5in footnote overprinted row 6."""
    snap = _rich_snap()
    snap["migration_readiness"] = (
        [{"group": "Group 1", "readiness": "NOT READY", "n_fail": 2, "n_warn": 3}]
        + [{"group": f"Group {i}", "readiness": "CAUTION", "n_fail": 0, "n_warn": 2}
           for i in range(2, 54)])
    snap["move_groups"] = [{"switches": [f"s{i}"]} for i in range(53)]
    snap["design_blueprint"] = {
        "summary": {"n_decisions": 1, "n_recommended": 1, "n_needs_requirement": 0, "n_critical": 1,
                    "headline": "1 recommended."},
        "tradeoff_scorecard": [], "coverage": {},
        "decisions": [{"id": "d1", "title": "T", "priority": "High", "status": "recommended",
                       "evidence": {"summary": "e"}, "principle": {"citation": "CCDE"}}],
        "target_state": {"wave_plan": {"n_waves": 9, "n_move_groups": 53, "largest_group": 251,
                                       "wave_cap": 40, "waves": [], "note": "n"}},
    }
    return snap


def _waves_slide(prs):
    for sl in prs.slides:
        if any(sh.has_text_frame and "PER-GROUP READINESS" in sh.text_frame.text for sh in sl.shapes):
            return sl
    raise AssertionError("migration-waves slide not found")


def test_deck_wave_footnote_below_list_by_construction(tmp_path):
    """QA row-11 BLOCK, defect 1 (baked-in shape overlap): the reconciling footnote was pinned at 6.5in
    while a 6-row readiness list runs to ~6.83in — the shapes overprinted and both were illegible. The
    footnote's top must now be COMPUTED from the rows actually rendered: strictly below every other shape
    on the slide, and fully on-slide. Bounding boxes asserted in EMU straight off the saved pptx."""
    out = tmp_path / "d.pptx"
    write_executive_deck_pptx(str(out), _m0_like_snap(), "Test fleet")
    sl = _waves_slide(Presentation(str(out)))
    foots = [sh for sh in sl.shapes
             if sh.has_text_frame and "Cutover Validation" in sh.text_frame.text]
    assert len(foots) == 1, "exactly one reconciling footnote expected on the waves slide"
    foot = foots[0]
    others_bottom = max(sh.top + sh.height for sh in sl.shapes if sh.shape_id != foot.shape_id)
    assert foot.top >= others_bottom, (
        f"footnote (top {foot.top} EMU) overlaps slide content (bottom {others_bottom} EMU)")
    assert foot.top + foot.height <= Presentation(str(out)).slide_height, "footnote runs off the slide"


def test_deck_wave_slide_discloses_hidden_groups_and_denominators(tmp_path):
    """QA row-11 BLOCK, defect 2 (readiness framing understates blocked scope ~7x): tiles
    '1 NOT READY / 52 CAUTION / 0 READY' sat beside '9 candidate waves' with a truncated list and NO cue
    that the rows shown are a subset of the 53-group estate. The slide must disclose the truncation
    ('+N more', computed from the data, never hardcoded) and carry the denominator on every readiness
    tile, and the wave tile must name the move-group scope it sequences."""
    out = tmp_path / "d.pptx"
    snap = _m0_like_snap()
    write_executive_deck_pptx(str(out), snap, "Test fleet")
    sl = _waves_slide(Presentation(str(out)))
    txt = "\n".join(sh.text_frame.text for sh in sl.shapes if sh.has_text_frame)
    n = len(snap["migration_readiness"])
    shown = txt.count("blocking ·")                   # one per rendered readiness row
    assert 0 < shown < n, "fixture must overflow the list for the cue to be exercised"
    assert f"+ {n - shown} more device group(s) not shown" in txt
    assert f"of {n} device group(s)" in txt           # tile denominators explicit
    assert "candidate waves · 53 move-group(s)" in txt


def test_deck_title_footer_carries_provenance(tmp_path):
    """P3-E3: the title footer must carry engine version + generation timestamp + snapshot stem (the
    explorer-ctx pattern, html.py:539-543) so a printed / forwarded deck stays traceable to the exact
    snapshot that produced it. The '_executive_deck' artifact suffix is stripped to the snapshot stem."""
    out = tmp_path / "Migration_Assessment_X_executive_deck.pptx"
    snap = _rich_snap()
    snap["script_version"] = "V3.23.0"
    snap["generated_at"] = "2026-07-11T10:27:30.004051"
    write_executive_deck_pptx(str(out), snap, "Test fleet")
    prs = Presentation(str(out))
    txt = "\n".join(sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame)
    assert "V3.23.0" in txt and "2026-07-11 10:27:30" in txt
    assert "snapshot Migration_Assessment_X" in txt
    assert "Migration_Assessment_X_executive_deck" not in txt


def test_deck_keystone_not_well_distributed_when_blast_radius_blind(tmp_path):
    """[audit-4 #8 false-health] the keystone slide printed green 'dependency is well distributed' whenever no
    record had stranded>0 -- including when failure_impact was NEVER computed (absent) or every record is
    INDETERMINATE (off-scan gateway). Absence/indeterminacy must read as a coverage gap, not a clean bill."""
    snap_absent = _rich_snap(); snap_absent.pop("failure_impact", None)
    write_executive_deck_pptx(str(tmp_path / "a.pptx"), snap_absent, "blind")
    ta = _deck_text(str(tmp_path / "a.pptx"))
    assert "well distributed" not in ta and ("not assess" in ta.lower() or "not computed" in ta.lower())
    snap_ind = _rich_snap()
    snap_ind["failure_impact"] = [{"host": "core1", "severity": "Info", "stranded": 0, "vlans_impacted": 0,
                                   "off_scan_gw_vlans": 5, "detail": "Blast radius INDETERMINATE — off-scan gateway"}]
    write_executive_deck_pptx(str(tmp_path / "b.pptx"), snap_ind, "ind")
    tb = _deck_text(str(tmp_path / "b.pptx"))
    assert "well distributed" not in tb and "indetermin" in tb.lower()
    snap_clean = _rich_snap()
    snap_clean["failure_impact"] = [{"host": "core1", "severity": "Low", "stranded": 0, "vlans_impacted": 0,
                                     "off_scan_gw_vlans": 0, "detail": "FHRP-covered"}]
    write_executive_deck_pptx(str(tmp_path / "c.pptx"), snap_clean, "clean")
    assert "well distributed" in _deck_text(str(tmp_path / "c.pptx"))
