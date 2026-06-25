"""Regression tests for the round-2 multi-agent code-review LOW/MEDIUM findings (engine side).

Each test names its finding id and pins the corrected behaviour so the bug class cannot silently
re-emerge. Webapp-side findings live under webapp/tests/.
"""
import logging

import pytest

from cisco_toolkit import analyze, portdb, ouidb
import cisco_toolkit.design_advisor as da


# ---------------------------------------------------------------- ANALY-04 ---
def test_blank_hostname_not_scored_so_reconcile_agrees():
    """ANALY-04: a blank/whitespace hostname (a malformed devices.json row that nonetheless collected) was
    scored by compute_health_scores but SKIPPED by compute_collection_completeness, so len(health_scores)
    exceeded the inventory count and ssot.reconcile false-fired an integrity DRIFT alarm. Both surfaces must
    now agree."""
    hs = analyze.compute_health_scores({"SW-A": {}, "SW-B": {}, "": {}}, [], [], [], [])
    assert sorted(r["switch"] for r in hs) == ["SW-A", "SW-B"]            # blank dropped, not scored
    cc = analyze.compute_collection_completeness(
        ["SW-A", "SW-B", ""],
        {"SW-A": {"show version": "x"}, "SW-B": {"show version": "x"}, "": {"show version": "x"}})
    assert len(hs) == cc["summary"]["inventory"]                          # surfaces agree -> no reconcile drift


# ---------------------------------------------------------------- ANALY-05 ---
def test_golden_drift_dossier_reads_na_when_no_baseline_derived():
    """ANALY-05: in majority mode with <3 comparable configs, compute_golden_drift derives NO baseline
    (summary.n_baseline == 0) yet still emits per_device rows with n_missing 0 / compliance 100. The dossier
    Golden-drift axis read those as 'ok / matches the config baseline' -- asserting conformance to a baseline
    that never existed (false-health). It must read 'na'."""
    gd = analyze.compute_golden_drift({"A": "hostname A\nip routing", "B": "hostname B"})
    assert gd["summary"]["n_baseline"] == 0
    dos = analyze.compute_device_dossiers(health_scores=[{"switch": "A"}, {"switch": "B"}], golden_drift=gd)
    states = [(e["state"], e["label"]) for pd in dos["per_device"]
              for e in pd["exposures"] if e["axis"] == "Golden drift"]
    assert states and all(s == "na" for s, _ in states), states
    assert all("no config baseline" in lbl for _, lbl in states)


# ---------------------------------------------------------------- LIFEC-04 ---
def test_near_ldos_band_compares_unrounded_delta():
    """LIFEC-04 (rounding half -- the other half, making the LDoS day itself 'Past', was REFUTED: LDoS = Last
    Day of Support, so support still exists ON that date; test_lifecycle_boundary_drift_guard locks the strict
    `>`). years_to_ldos was ROUNDED before the <= 1.0 band test, so a device genuinely ~1.05 yr from LDoS
    rounded to 1.0 and fell into Near-LDoS (the band ~18 days too wide). The band must use the UNROUNDED
    delta. (Catalyst 3850 LDoS = 2026-10-31.)"""
    band = lambda asof: analyze.compute_lifecycle_risk(
        {"sw1": {"model": "WS-C3850-48P"}}, asof=asof)["per_device"][0]["band"]
    assert band("2025-10-13") != "Near-LDoS"     # ~1.05 yr out: was Near via rounding (1.0), now Past-EoS
    assert band("2026-06-01") == "Near-LDoS"      # genuinely within 1 yr still Near
    assert band("2026-10-31") == "Near-LDoS"      # LDoS day: last supported day -> Near, not yet Past (strict >)


# ---------------------------------------------------------------- LIFEC-03 ---
def test_offline_registries_warn_on_load_failure_then_degrade(caplog):
    """LIFEC-03: portdb/ouidb._registry() swallowed ANY load error with a bare `except Exception: pass` and,
    being lru_cached, memoised the empty result for the whole process with NO signal -- a missing/corrupt
    shipped pack silently disabled the entire L4-service/multicast and MAC->vendor axes. A load failure must
    now emit a WARNING (and still degrade tolerantly to empty, never crash)."""
    assert portdb.service_for_port(179, "tcp") is not None                # sanity: the real pack loads
    p_orig, o_orig = portdb._DATA, ouidb._DATA
    try:
        portdb._DATA = "/nonexistent.gz"; portdb._registry.cache_clear()
        ouidb._DATA = "/nonexistent.gz"; ouidb._registry.cache_clear()
        with caplog.at_level(logging.WARNING):
            ports, mcast = portdb._registry()
            tables = ouidb._registry()
        assert ports == {} and mcast == []                               # tolerant degrade, no crash
        assert all(v == {} for v in tables.values())
        msgs = " ".join(r.getMessage() for r in caplog.records)
        assert "portdb: could not load" in msgs                          # surfaced, not silent
        assert "ouidb: could not load" in msgs
    finally:                                                              # MUST restore or poison the suite
        portdb._DATA = p_orig; portdb._registry.cache_clear()
        ouidb._DATA = o_orig; ouidb._registry.cache_clear(); ouidb.vendor_for_mac.cache_clear()
    assert portdb.service_for_port(179, "tcp") is not None                # restored cleanly


# ---------------------------------------------------------------- DETEC-03 ---
def test_fhrp_no_preempt_fires_only_on_raised_priority_not_default():
    """DETEC-03: the signal flagged ANY active gateway with preempt off. But HSRP/VRRP/GLBP preemption is OFF
    BY DEFAULT (and parse_hsrp_detail defaults preempt=False), so every textbook default active group cried
    wolf. Only a gateway whose operator RAISED its configured priority above 100 deliberately wants to be
    primary -- without preempt that intended primary won't reclaim after a failover. A default-priority (100)
    active gateway with no preempt must stay silent."""
    default = {"fhrp_detail": {"d1": [{"ifname": "Vlan10", "group": "10", "state": "Active",
                                       "preempt": False, "priority": 100, "cfg_priority": 100,
                                       "track": [{"obj": "1", "decrement": 10}]}]}}
    raised = {"fhrp_detail": {"d2": [{"ifname": "Vlan20", "group": "20", "state": "Active",
                                      "preempt": False, "priority": 110, "cfg_priority": 110,
                                      "track": [{"obj": "1", "decrement": 10}]}]}}
    assert da._signals(default)["fhrp_no_preempt"] == []                  # benign default: silent
    assert da._signals(raised)["fhrp_no_preempt"] == ["d2 Vlan20 grp 20"]  # raised + no preempt: fires
    assert da._d_fhrp_resilience(default, da._signals(default)) is None    # no other arm fires either


# ---------------------------------------------------------------- DETEC-04 ---
def test_sdwan_control_shortfall_counted_once_per_device():
    """DETEC-04: vManage repeats the DEVICE-level expected/actual control-connection count on every per-peer
    row, so a single device short by one made every remaining UP row also satisfy actual<expected -- counting
    the same shortfall N times. It must be counted ONCE per device; a genuinely-down peer is still per-row."""
    rows = [{"system_ip": "10.0.0.13", "host_name": "BR13", "peer_type": "vsmart", "state": "up",
             "expected": 4, "actual": 3} for _ in range(3)]
    rows.append({"system_ip": "10.0.0.13", "host_name": "BR13", "peer_type": "vbond", "state": "down",
                 "expected": 4, "actual": 3})
    cc = da._signals({"sdwan": {"mgr1": {"control_connections": rows}}})["sdwan_control_down"]
    assert len(cc) == 2, cc                                               # was 4 (one per row); now 1 short + 1 down
    assert sum("control connections 3/4" in x for x in cc) == 1          # the shortfall: exactly once
    assert sum("(down)" in x for x in cc) == 1                            # the down peer: still surfaced


# ---------------------------------------------------------------- BUILD-02 ---
def test_build_cloud_coverage_honest_on_parser_crash(tmp_path):
    """BUILD-02: build_cloud's not-observed contract hinged on `sgs is None`, but _safe_parse returns {} (a
    truthy non-list) when the parser RAISES -- so a crash produced {security_groups: {}}, read downstream as
    'cloud observed, nothing world-open' (false-health). A non-list (crash {} or no-export None) must map to
    not-observed {}; only a real list is a result."""
    import cisco_toolkit.build as b
    f = tmp_path / "sg.json"; f.write_text('{"SecurityGroups":[]}')
    orig = b.parse_aws_security_groups
    try:
        b.parse_aws_security_groups = lambda *a: (_ for _ in ()).throw(RecursionError())
        assert b.build_cloud({"aws ec2 describe-security-groups": str(f)}) == {}        # crash -> not observed
    finally:
        b.parse_aws_security_groups = orig
    assert b.build_cloud({}) == {}                                                      # no export -> not observed
    assert b.build_cloud({"aws ec2 describe-security-groups": str(f)}) == {"security_groups": []}  # observed clean


# ---------------------------------------------------------------- BUILD-03 ---
def test_archreview_sheet_distinguishes_crash_from_legit_empty():
    """BUILD-03: a crashed/absent architecture review (_run_phase _default) rendered as a clean 'grade N/A · 0
    issues' scorecard -- false-health. The sheet must render 'unavailable' for the crash sentinel / bare {},
    while a legitimately-empty review (which still carries a summary, n_not_assessable counting un-evidenced
    checks) keeps its honest grade row."""
    from openpyxl import Workbook
    from cisco_toolkit.archreview import compute_architecture_review
    from cisco_toolkit.excel import ARCHREVIEW_SHEET_NAME as N, write_architecture_review_sheet as W
    wb_legit = Workbook(); W(wb_legit, compute_architecture_review({}))
    wb_crash = Workbook(); W(wb_crash, {})
    wb_sentinel = Workbook(); W(wb_sentinel, {"_unavailable": True})
    assert "unavailable" not in wb_legit[N].cell(2, 1).value.lower()      # legit-empty keeps its honest grade row
    assert "unavailable" in wb_crash[N].cell(2, 1).value.lower()          # bare {} crash -> disclosed
    assert "unavailable" in wb_sentinel[N].cell(2, 1).value.lower()       # sentinel -> disclosed


# ------------------------------------------------------- CROSS-01 / CROSS-02 ---
def test_crd_and_runbook_device_count_is_canonical_not_raw_len(tmp_path):
    """CROSS-01/02: crd._evidence_facts and the runbook 'Devices in scope' row recomputed n_devices from raw
    len(devices) while their sibling n_vlans/n_endpoints read executive_brief.scale canonical-first -- the exact
    device-count drift seam ssot.py exists to eliminate. Both must read the canonical count (the explorer/deck/
    HLD value), falling back to len() only pre-brief."""
    from cisco_toolkit import crd
    from cisco_toolkit.runbook import write_runbook_docx
    from docx import Document
    snap = {"executive_brief": {"scale": {"n_devices": 2, "n_endpoints": 10, "n_vlans": 5}},
            "devices": {"a": {}, "b": {}, "c": {}, "d": {}}, "health_scores": [], "punchlist": []}
    assert crd._evidence_facts(snap)["n_devices"] == 2                    # CROSS-01: canonical 2, not len()=4
    p = tmp_path / "rb.docx"; write_runbook_docx(str(p), snap, "Test")
    cells = [c.text for t in Document(str(p)).tables for r in t.rows for c in r.cells]
    i = next(k for k, c in enumerate(cells) if "Devices in scope" in c)
    assert cells[i + 1] == "2"                                            # CROSS-02: canonical 2, not 4


def test_deck_title_scale_renders_dash_not_none_for_null_scale(tmp_path):
    """DECK_-02: the deck title slide built the scale line with scale.get('n_devices', '—'); dict.get's default
    only fires when the KEY is ABSENT, so a present-but-null value (uploaded / partially-computed snapshot)
    rendered the literal 'None devices · None endpoints' on the marquee client slide. Coerce per value -> em-dash."""
    from cisco_toolkit.deck import write_executive_deck_pptx
    from pptx import Presentation
    snap = {"executive_brief": {"scale": {"n_devices": None, "n_endpoints": None, "n_vlans": None},
                                "posture": {}, "posture_statement": "Review.", "axes": [], "top_gating": []},
            "devices": {"core1": {}}, "health_scores": [{"switch": "core1", "score": 50, "band": "Fair"}],
            "punchlist": [], "failure_impact": [], "migration_readiness": [], "move_groups": [], "lifecycle_risk": {}}
    p = tmp_path / "deck.pptx"; write_executive_deck_pptx(str(p), snap, "Test")
    blob = " ".join(sh.text_frame.text for sl in Presentation(str(p)).slides
                    for sh in sl.shapes if sh.has_text_frame)
    assert "None devices" not in blob and "None endpoints" not in blob
    assert "— devices" in blob                                           # em-dash, not the literal 'None'


# --- shared real 'show vpc' capture: healthy domain, 1 down leg (Po20), 1 UP-but-inconsistent leg (Po30) ---
_VPC_SHOW = (
    "vPC domain id                     : 10\n"
    "Peer status                       : peer adjacency formed ok\n"
    "vPC keep-alive status             : peer is alive\n"
    "Configuration consistency status  : success\n"
    "vPC role                          : primary\n"
    "Number of vPCs configured         : 3\n\n"
    "vPC Peer-link status\n"
    "---------------------------------------------------------------------\n"
    "id    Port   Status Active vlans\n"
    "--    ----   ------ --------------\n"
    "1     Po1    up     1,10,20,30\n\n"
    "vPC status\n"
    "----------------------------------------------------------------------------\n"
    "Id    Port          Status Consistency Reason                Active vlans\n"
    "--    ------------  ------ ----------- ------                ------------\n"
    "10    Po10          up     success     success               1,10,20\n"
    "20    Po20          down   failed      down                  1,20\n"
    "30    Po30          up     failed      compat-failed         1,30\n")


# ---------------------------------------------------------------- TEST-01 ---
def test_vpc_detector_end_to_end_from_real_parser():
    """TEST-01: _d_vpc_health had no end-to-end test running the REAL parse_vpc into the detector -- the
    parser-shape <-> detector-read agreement for vPC was unverified (no 'show vpc' capture in synthetic_fixtures
    either). A genuinely-degraded vPC (a down member leg + an UP-but-inconsistent leg, on a HEALTHY domain) must
    fire; an all-up/success fabric stays silent."""
    from cisco_toolkit import parse
    snap = {"vpc": {"nx-core1": parse.parse_vpc(_VPC_SHOW)}}
    sig = da._signals(snap)
    assert sig["vpc_legs"] == 3 and sig["vpc_legs_down"] == 1 and sig["vpc_legs_incons"] == 1
    assert sig["vpc_dom_bad"] == 0                            # domain itself healthy (peer ok / cons success)
    assert da._d_vpc_health(snap, sig) is not None           # fires on the degraded legs
    clean = (_VPC_SHOW.replace("down   failed      down", "up     success     success")
                      .replace("up     failed      compat-failed", "up     success     success"))
    csnap = {"vpc": {"nx-core1": parse.parse_vpc(clean)}}
    csig = da._signals(csnap)
    assert csig["vpc_unhealthy"] == 0 and da._d_vpc_health(csnap, csig) is None


# ---------------------------------------------------------------- TEST-02 ---
def test_port_security_detail_nxos_headerless_form():
    """TEST-02: the NX-OS 'show port-security interface ethernet 1/1' form carries the interface in the COMMAND,
    not echoed in the body, so a header-less Secure-shutdown block parsed to {} -- and the test bank (IOS-only +
    {}-on-noise) made that blind spot look correct. Without context the block stays {} (a documented residual,
    NOT silently 'healthy'); a caller that knows the interface supplies it via default_ifname and the body fields
    (identical IOS/NX-OS) then parse the err-disabled state."""
    from cisco_toolkit import parse
    nxos = ("Port Security                 : Enabled\n"
            "Port Status                   : Secure-shutdown\n"
            "Violation Mode                : Shutdown\n"
            "Security Violation Count      : 7\n")
    assert parse.parse_port_security_detail(nxos) == {}                  # header-less, no ctx -> not silently healthy
    r = parse.parse_port_security_detail(nxos, default_ifname="Ethernet1/1")
    assert r["Eth1/1"]["port_status"] == "secure-shutdown"              # the live-outage state surfaces
    assert r["Eth1/1"]["enabled"] is True and r["Eth1/1"]["violation_count"] == 7
    assert "Gi0/1" in parse.parse_port_security_detail("Port: Gi0/1\nPort Security : Enabled\n")   # IOS unchanged


# ---------------------------------------------------------------- TEST-03 ---
def test_parser_to_detector_contract_via_real_build(tmp_path):
    """TEST-03: no test asserted, AS A CATEGORY, that a detector reads the REAL build/parser output shape -- the
    hand-built dict fixtures in test_design_blueprint could drift from what build_X actually emits and silently
    zero a detector. This runs _signals(build_X(real_text)) for switch-native axes so a parser-shape drift fails
    fast. (vpc was the proven gap in TEST-01; fhrp_detail adds a second axis.)"""
    from cisco_toolkit import build

    def _c2f(cmd, text):
        p = tmp_path / (cmd.replace(" ", "_") + ".txt"); p.write_text(text, encoding="utf-8"); return {cmd: str(p)}

    vsnap = {"vpc": {"core1": build.build_vpc(_c2f("show vpc", _VPC_SHOW))}}
    vsig = da._signals(vsnap)
    assert vsig["vpc_legs"] == 3 and vsig["vpc_legs_down"] == 1          # build_vpc -> vpc signal shape agrees

    standby = ("Vlan10 - Group 10 (version 2)\n  State is Active\n  Preemption disabled\n"
               "  Priority 130 (configured 130)\n")                       # raised priority + no preempt (DETEC-03)
    fsnap = {"fhrp_detail": {"core1": build.build_fhrp_detail(_c2f("show standby all", standby))}}
    fsig = da._signals(fsnap)
    assert fsig["fhrp_no_preempt"] == ["core1 Vlan10 grp 10"]            # build_fhrp_detail -> fhrp signal shape agrees
