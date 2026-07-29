"""Round-6 silent-truncation sweep: every `[:N]` DISPLAY cap in the MOP (`cisco_toolkit/mop.py`) and
the executive deck (`cisco_toolkit/deck.py`) must DISCLOSE what it dropped.

The house rule (stated in `excel.py :: _xls_cell_value`) is that a cap carries a trailing "(+N)" marker
or an "…and N further …" sentence. Round 4 fixed the §x.2 blocker lists in this same writer, where a
capped list sat under a PRECONDITION GATE that stated the full attributed counts — an engineer worked
the 15 printed items, ticked "blockers cleared", and would have opened a production window on 70
uncleared Critical/High blockers. These tests pin the remaining sites of that shape: a cap whose full
count is stated elsewhere in the same document, or which a gate/checklist depends on.

Every cap pinned here was MEASURED as reachable on the real producer artifacts before being fixed
(repo-root `Migration_Assessment_AUTOFILLED_20260613_063201.snapshot.json`, 303 devices / 9 waves):
9 fleet-wide gating items vs a 6 cap; up to 35 management ports per wave vs 6; up to 40 stub-bearing
devices per wave vs 2 and up to 95 endpoint ports on one device vs 6; 193 switches stranding endpoints
vs a top-5 slide; 30 recommended design decisions vs 5 under a headline that states "30 recommended";
6 trade-off axes at 0-1/4 vs 5 chips; and asset verdicts of 186-217 chars against a bare [:160] that
ended "…in the device's own logs; removal strand".
"""
import re

import pytest

docx = pytest.importorskip("docx")          # optional dependency, exactly as test_mop.py skips
pptx = pytest.importorskip("pptx")          # optional dependency, exactly as test_deck.py skips

from docx import Document                                       # noqa: E402
from pptx import Presentation                                   # noqa: E402

from cisco_toolkit.deck import write_executive_deck_pptx        # noqa: E402
from cisco_toolkit.mop import write_mop_docx                    # noqa: E402


# --------------------------------------------------------------------------- helpers

def _mop_snap(switches, **over):
    """A single-wave MOP snapshot; `switches` are the wave's devices."""
    snap = {
        "script_version": "V3.23.0",
        "devices": {h: {"platform": "ios"} for h in switches},
        "move_groups": [{"switches": list(switches), "endpoints": 10}],
        "migration_readiness": [{"group": "Group 1", "switches": list(switches), "endpoints": 10,
                                 "readiness": "CAUTION", "n_fail": 0, "n_warn": 1, "checks": []}],
        "executive_brief": {"top_gating": []},
    }
    snap.update(over)
    return snap


def _paras(doc):
    return [p.text for p in doc.paragraphs]


def _row_values(doc, label, col=2):
    """The `col`-th cell of every table row whose FIRST cell is exactly `label` (test_mop.py's
    mutation-proved lookup: a whole-document `in text` assertion cannot tell a rendered ROW from the
    boilerplate paragraph that merely explains the convention — look the row up)."""
    out = []
    for t in doc.tables:
        for row in t.rows:
            cells = [c.text.strip() for c in row.cells]
            if cells and cells[0] == label and len(cells) > col:
                out.append(cells[col])
    return out


def _deck_snap(**over):
    snap = {
        "executive_brief": {"scale": {"n_devices": 9, "n_endpoints": 100, "n_vlans": 5},
                            "posture": {"avg_health": 50, "n_critical": 1},
                            "posture_statement": "Migration posture: unit fixture.",
                            "axes": [], "top_gating": []},
        "health_scores": [{"switch": "sw1", "score": 50, "band": "Fair"}],
        "punchlist": [], "failure_impact": [], "migration_readiness": [], "move_groups": [],
        "lifecycle_risk": {"summary": {}},
    }
    snap.update(over)
    return snap


def _deck_text(path):
    p = Presentation(path)
    return "\n".join(sh.text_frame.text for sl in p.slides for sh in sl.shapes if sh.has_text_frame)


# --------------------------------------------------------------------------- MOP §2 global prerequisites

def test_mop_global_prerequisites_disclose_dropped_gating_items(tmp_path):
    """§2 "Global Prerequisites (before any window)" printed the first SIX fleet-wide gating items and
    dropped the rest with no cue, while its own lead-in makes the list a gate: "Complete these once,
    before the first wave … cutting over before they are resolved or risk-accepted carries the
    documented risk." The real fleet publishes NINE; the deck already discloses the same list as
    "5 of 9 shown", so a silent MOP was the set-wide outlier. Hidden items 7-9 there were a MAC clash,
    3788 critical log events across 253 devices, and 109 exposed advisory surfaces."""
    snap = _mop_snap(["sw1", "sw2"])
    snap["executive_brief"] = {"top_gating": [f"gating item {i}" for i in range(1, 10)]}
    out = str(tmp_path / "mop_gating.docx")
    write_mop_docx(out, snap, "Unit Test Fleet")
    paras = _paras(Document(out))
    body = "\n".join(paras)

    assert "gating item 6" in body and "gating item 7" not in body, "cap moved; retune this test"
    cue = [t for t in paras if "further fleet-wide gating item" in t]
    assert cue, ("§2 dropped 3 of 9 fleet-wide gating items with no disclosure — an engineer who "
                 f"clears the 6 printed believes the prerequisite is met. Paragraphs: {paras[:20]}")
    shown, hidden, total = 6, int(re.search(r"…and (\d+) further", cue[0]).group(1)), \
        int(re.search(r"\((\d+) in total\)", cue[0]).group(1))
    assert shown + hidden == total == 9, f"shown+hidden must reconcile to the real total: {cue[0]!r}"
    assert "not only the 6 shown" in cue[0], cue[0]


def test_mop_global_prerequisites_stay_silent_when_nothing_is_dropped(tmp_path):
    """Negative control: a brief carrying <= 6 gating items must render NO overflow sentence (a cue
    that always fires is noise, and would tell an engineer to hunt for items that do not exist)."""
    snap = _mop_snap(["sw1", "sw2"])
    snap["executive_brief"] = {"top_gating": [f"gating item {i}" for i in range(1, 7)]}
    out = str(tmp_path / "mop_gating_exact.docx")
    write_mop_docx(out, snap, "Unit Test Fleet")
    assert "further fleet-wide gating item" not in "\n".join(_paras(Document(out)))


# --------------------------------------------------------------------------- MOP OOB precondition row

def test_mop_oob_precondition_names_the_management_ports_it_dropped(tmp_path):
    """The pre-implementation checklist's OOB row is the EVIDENCE gate for "Out-of-band / console
    access to every device in scope". It listed six management ports and appended a bare "…" — which
    marks THAT something was dropped but never WHAT: on the real fleet one wave observes 35 ports and
    printed 6, so the engineer confirms 6 consoles and ticks a gate covering 35. Losing OOB to an
    unlisted device mid-window is losing the rollback path."""
    hosts = [f"sw{i:02d}" for i in range(1, 10)]                 # 9 devices, each with a mgmt port
    snap = _mop_snap(hosts)
    snap["interfaces"] = {h: {"mgmt0": {"mgmt_ip": f"10.99.0.{i}"}} for i, h in enumerate(hosts, 1)}
    out = str(tmp_path / "mop_oob.docx")
    write_mop_docx(out, snap, "Unit Test Fleet")
    rows = _row_values(Document(out), "Out-of-band / console access to every device in scope")
    assert len(rows) == 1, f"expected one OOB precondition row per wave, got {rows}"
    row = rows[0]

    assert row.count(";") == 5 and "sw06 mgmt0" in row, f"cap moved; retune this test: {row!r}"
    m = re.search(r"…and (\d+) further management-port\(s\) not listed here \((\d+) observed in total\)", row)
    assert m, (f"the OOB evidence gate hid {len(hosts) - 6} of {len(hosts)} management ports behind an "
               f"anonymous ellipsis: {row!r}")
    assert 6 + int(m.group(1)) == int(m.group(2)) == len(hosts), row
    assert "confirm each is reachable before the window" in row


# --------------------------------------------------------------------------- MOP §x.4 staged config

def _iface_snap(spec):
    """spec: {host: n_endpoint_ports} -> a snapshot whose §x.4 builds one stub block per host."""
    hosts = sorted(spec)
    snap = _mop_snap(hosts)
    snap["interfaces"] = {
        h: {f"Gi1/0/{p}": {"switchport_mode": "Access", "vlan": "20",
                           "end_host_mac": f"bbbb.{i:04x}.{p:04x}", "description": f"EP-{h}-{p}"}
            for p in range(1, spec[h] + 1)}
        for i, h in enumerate(hosts)
    }
    return snap


def test_mop_staged_config_discloses_the_devices_it_did_not_stage(tmp_path):
    """§x.4 "Port / interface mapping & staged configuration" stubs the first TWO devices only, and
    §x.5 then instructs "Apply the staged target configuration … per the §x.4 port mapping". On the
    real fleet seven of nine waves carry more than two stub-bearing devices (up to 40), so the
    document silently presented 2-of-40 devices' configuration as the wave's staged configuration."""
    snap = _iface_snap({"aa-sw1": 2, "bb-sw2": 2, "cc-sw3": 2, "dd-sw4": 2, "ee-sw5": 2})
    out = str(tmp_path / "mop_stub_hosts.docx")
    write_mop_docx(out, snap, "Unit Test Fleet")
    paras = _paras(Document(out))
    blocks = [t for t in paras if "staged endpoint-port config (evidence-derived)" in t]
    assert len(blocks) == 2 and "cc-sw3" not in "\n".join(blocks), "cap moved; retune this test"

    cue = [t for t in paras if "further device(s) in this wave carry endpoint ports" in t]
    assert cue, (f"3 of the wave's 5 stub-bearing devices were dropped with no disclosure; "
                 f"paragraphs: {[t[:70] for t in paras if 'staged' in t]}")
    assert re.search(r"…and 3 further device\(s\)", cue[0]), cue[0]
    assert "(5 in total)" in cue[0] and "only the first 2 are stubbed above" in cue[0], cue[0]
    assert "NOT this wave's complete staged configuration" in cue[0], cue[0]


def test_mop_staged_config_discloses_the_ports_it_did_not_stage(tmp_path):
    """The second, nested cap in the same block: only the first SIX endpoint ports of a stubbed device
    get an interface stanza. The real fleet has a device with 95 endpoint ports in a wave, so the
    'staged configuration' for it covered 6 — with nothing saying so, inside a Consolas block that
    reads as a complete, paste-ready config."""
    snap = _iface_snap({"aa-sw1": 9, "bb-sw2": 2})
    out = str(tmp_path / "mop_stub_ports.docx")
    write_mop_docx(out, snap, "Unit Test Fleet")
    block = next(t for t in _paras(Document(out)) if t.startswith("! aa-sw1 —"))
    stanzas = [ln for ln in block.split("\n") if ln.startswith("interface <target port")]
    assert len(stanzas) == 6, f"cap moved; retune this test: {len(stanzas)} stanzas"

    tail = [ln for ln in block.split("\n") if "further endpoint port(s)" in ln]
    assert tail, f"3 of aa-sw1's 9 endpoint ports were dropped from the staged config silently:\n{block}"
    assert tail[0].startswith("!"), f"the disclosure must stay a config COMMENT: {tail[0]!r}"
    assert "…and 3 further endpoint port(s) on aa-sw1" in tail[0] and "(9 in total)" in tail[0], tail[0]
    # the un-truncated device must NOT carry a cue (shown + hidden reconciles per device)
    b2 = next(t for t in _paras(Document(out)) if t.startswith("! bb-sw2 —"))
    assert "further endpoint port(s)" not in b2, b2


# --------------------------------------------------------------------------- deck: keystone switches

def test_deck_keystone_slide_discloses_the_switches_below_the_top_five(tmp_path):
    """Slide "The switches the fleet depends on" ranks by blast radius and says "Sequence and protect
    these first" — with no cue that it shows five. Every sibling list on this deck breadcrumbs its
    overflow. On the real fleet 193 switches strand endpoints and the top TEN are TIED at 2348, so the
    five unmarked rows hid five more switches of identical blast radius."""
    fi = [{"host": f"sw{i:02d}", "severity": "High", "stranded": 100 - i, "vlans_impacted": 2,
           "hard": 1, "detail": f"VLAN {i} hard partition"} for i in range(1, 10)]
    out = str(tmp_path / "deck_keystones.pptx")
    write_executive_deck_pptx(out, _deck_snap(failure_impact=fi), "Unit fleet")
    txt = _deck_text(out)

    assert "sw05" in txt and "sw06" not in txt, "cap moved; retune this test"
    assert "Top 5 of 9 switch(es) that strand endpoints shown" in txt, (
        "the keystone slide showed 5 of 9 blast-radius switches with no disclosure — a steering "
        "committee reads the five as the whole keystone set")
    assert "Failure-Impact sheet" in txt


def test_deck_keystone_slide_stays_silent_when_nothing_is_dropped(tmp_path):
    """Negative control: <= 5 stranding switches must render no 'Top N of M' clause."""
    fi = [{"host": f"sw{i}", "stranded": 10 * i, "vlans_impacted": 1, "detail": "d"} for i in range(1, 4)]
    out = str(tmp_path / "deck_keystones_small.pptx")
    write_executive_deck_pptx(out, _deck_snap(failure_impact=fi), "Unit fleet")
    assert "switch(es) that strand endpoints shown" not in _deck_text(out)


# --------------------------------------------------------------------------- deck: design decisions

def _bp_snap(n_rec, n_weak=0):
    decisions = [{"id": f"d{i}", "title": f"Decision {i:02d}", "priority": "Critical",
                  "status": "recommended", "evidence": {"summary": f"evidence {i}"}}
                 for i in range(1, n_rec + 1)]
    axes = [{"axis": f"axis{i}", "score": 0 if i <= n_weak else 3} for i in range(1, 11)]
    return _deck_snap(design_blueprint={
        "summary": {"n_decisions": n_rec, "n_recommended": n_rec, "n_critical": n_rec,
                    "n_needs_requirement": 0},
        "decisions": decisions, "tradeoff_scorecard": axes, "coverage": {}})


def test_deck_design_slide_discloses_the_recommendations_it_dropped(tmp_path):
    """The target-state slide's headline states the FULL population ("39 design decision(s) · 30
    recommended") and then listed five rows with no cue — the "full count stated elsewhere in the same
    document" shape. The rows run to y≈7.47in, so the disclosure reserves the last row (the computed
    layout rule slide 6 already uses) rather than overflowing the 7.5in slide."""
    out = str(tmp_path / "deck_decisions.pptx")
    write_executive_deck_pptx(out, _bp_snap(30), "Unit fleet")
    txt = _deck_text(out)
    p = Presentation(out)

    assert "30 recommended" in txt, "the headline must still state the full population"
    cue = [ln for ln in txt.split("\n") if "more recommended decision(s)" in ln]
    assert cue, f"25 of 30 recommended decisions were dropped with no cue:\n{txt[-900:]}"
    shown = sum(1 for ln in txt.split("\n") if ln.startswith("Decision "))
    hidden = int(re.search(r"\+ (\d+) more recommended decision", cue[0]).group(1))
    assert shown + hidden == 30, f"shown ({shown}) + hidden ({hidden}) must reconcile to 30: {cue[0]!r}"
    assert "design blueprint (HLD/LLD)" in cue[0], cue[0]
    # the reserved cue line must not push any shape off the 7.5in slide
    bottoms = [(sh.top + sh.height) / 914400 for sl in p.slides for sh in sl.shapes]
    assert max(bottoms) <= p.slide_height / 914400 + 1e-6, f"a shape overflows the slide: {max(bottoms)}"


def test_deck_design_slide_keeps_five_rows_when_nothing_is_dropped(tmp_path):
    """Negative control: with exactly 5 recommendations all five render and no line is reserved."""
    out = str(tmp_path / "deck_decisions_five.pptx")
    write_executive_deck_pptx(out, _bp_snap(5), "Unit fleet")
    txt = _deck_text(out)
    assert sum(1 for ln in txt.split("\n") if ln.startswith("Decision ")) == 5
    assert "more recommended decision(s)" not in txt


def test_deck_weakest_tradeoff_axes_disclose_the_axis_they_dropped(tmp_path):
    """A label promising "WEAKEST TRADE-OFF AXES (0-1 / 4)" must not silently drop one: the real fleet
    scores SIX axes at 0-1/4 and the slide holds five chips, so the sixth weak axis vanished under a
    heading that claims to be the weak set."""
    out = str(tmp_path / "deck_axes.pptx")
    write_executive_deck_pptx(out, _bp_snap(3, n_weak=6), "Unit fleet")
    txt = _deck_text(out)
    label = next(ln for ln in txt.split("\n") if ln.startswith("WEAKEST TRADE-OFF AXES"))
    assert "5 OF 6 SHOWN" in label, f"the 6th weak axis was dropped silently: {label!r}"
    assert sum(1 for ln in txt.split("\n") if ln.startswith("axis")) == 5, txt


def test_deck_weakest_tradeoff_axes_label_clean_when_all_shown(tmp_path):
    """Negative control: <= 5 weak axes leaves the label untouched."""
    out = str(tmp_path / "deck_axes_small.pptx")
    write_executive_deck_pptx(out, _bp_snap(3, n_weak=4), "Unit fleet")
    label = next(ln for ln in _deck_text(out).split("\n") if ln.startswith("WEAKEST TRADE-OFF AXES"))
    assert "SHOWN" not in label, label


# --------------------------------------------------------------------------- deck: asset verdict

def test_deck_asset_verdict_truncation_is_marked(tmp_path):
    """The riskiest-assets slide cut each engineer verdict with a bare `[:160]` — the exact class
    `_ellip` exists for ("a bare slice ends mid-word with no cue that text was dropped"). On the real
    fleet all five rendered verdicts run 186-217 chars, and the top row ended "…in the device's own
    logs; removal strand", dropping "s 286 endpoint(s) across 3 VLAN(s)." — the verdict's quantified
    impact — with nothing to show the sentence had been cut."""
    verdict = ("Stabilize or replace before migration — " + "compound exposure; " * 7
               + "removal strands 286 endpoint(s) across 3 VLAN(s).")
    assert len(verdict) > 160, "fixture must exceed the display cap"
    snap = _deck_snap(device_dossiers={
        "summary": {"bands": {"Severe": 1}},
        "per_device": [{"host": "sw1", "risk_band": "Severe", "risk_index": 90,
                        "compound": [{"code": "CR-01"}], "verdict": verdict}]})
    out = str(tmp_path / "deck_verdict.pptx")
    write_executive_deck_pptx(out, snap, "Unit fleet")
    line = next(ln for ln in _deck_text(out).split("\n") if ln.startswith("Stabilize or replace"))

    assert line != verdict[:160], "the verdict was cut mid-word with no cue that anything was dropped"
    assert line.endswith("…"), f"a capped verdict must mark the cut: {line!r}"
    assert len(line) <= 161, f"the cap itself must still hold: {len(line)}"


def test_deck_short_asset_verdict_is_rendered_whole(tmp_path):
    """Negative control: a verdict inside the cap renders verbatim, with no ellipsis added."""
    verdict = "Stabilize or replace before migration — hardware Past LDoS."
    snap = _deck_snap(device_dossiers={
        "summary": {"bands": {"Severe": 1}},
        "per_device": [{"host": "sw1", "risk_band": "Severe", "risk_index": 90,
                        "compound": [], "verdict": verdict}]})
    out = str(tmp_path / "deck_verdict_short.pptx")
    write_executive_deck_pptx(out, snap, "Unit fleet")
    assert verdict in _deck_text(out)
