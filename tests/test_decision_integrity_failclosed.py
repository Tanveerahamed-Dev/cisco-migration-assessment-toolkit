"""Fail-closed decision/output regressions from the whole-repository 2026-07 review."""

from copy import deepcopy

import pytest

from cisco_toolkit import assertions, cutover_sim, path_assertions, precert
from cisco_toolkit.html import (
    compute_campaign_trend,
    compute_snapshot_delta,
    write_diff_workbook,
)


def _comparable_snapshot(**extra):
    snap = {
        "devices": {},
        "interfaces": {},
        "health_scores": [],
        "punchlist": [],
        "migration_readiness": [],
        "lifecycle_risk": {"summary": {}},
    }
    snap.update(extra)
    return snap


def test_failed_empty_producer_never_resolves_findings_or_certifies_clean():
    finding = {
        "severity": "Critical",
        "category": "L3",
        "title": "gateway regression",
        "devices": ["core1"],
    }
    before = _comparable_snapshot(punchlist=[finding])
    after = _comparable_snapshot(
        punchlist=[],
        assessment_integrity={"failed_phases": ["Migration Punch-List"]},
    )

    delta = compute_snapshot_delta(before, after)

    assert delta["verdict"] == "INDETERMINATE"
    assert delta["findings"]["n_resolved"] == 0
    assert delta["integrity"]["findings_comparable"] is False
    assert "Migration Punch-List" in " ".join(delta["integrity"]["failures"])


def test_schema_mismatch_dominates_an_apparent_regression():
    before = _comparable_snapshot(
        health_scores=[{"switch": "sw1", "band": "Good", "score": 80}],
    )
    after = _comparable_snapshot(
        health_scores=[{"switch": "sw1", "band": "Poor", "score": 40}],
    )

    delta = compute_snapshot_delta(
        before,
        after,
        schema_status={"status": "mismatch", "message": "V1 cannot be certified against V2"},
    )

    assert delta["health"]["n_regressed"] == 1
    assert delta["verdict"] == "INDETERMINATE"
    assert "apparent health-band regression" in delta["verdict_note"]
    assert "schema compatibility mismatch" in " ".join(delta["integrity"]["failures"])


def test_failed_campaign_collection_cannot_manufacture_an_improving_trajectory():
    before = _comparable_snapshot(
        health_scores=[{"switch": "sw1", "band": "Critical", "score": 20}],
        punchlist=[{"severity": "Critical", "category": "L3", "title": "x", "devices": ["sw1"]}],
    )
    after = _comparable_snapshot(
        assessment_integrity={
            "failed_phases": ["Health Scores", "Migration Punch-List", "Lifecycle risk"]
        }
    )

    trend = compute_campaign_trend([before, after])

    assert trend["verdict"] == "INDETERMINATE"
    assert trend["integrity"]["ok"] is False
    assert "not an improvement claim" in trend["verdict_note"]


def test_clean_is_explicitly_delta_scoped_and_workbook_reconciles_certificate(tmp_path):
    snap = _comparable_snapshot()
    delta = compute_snapshot_delta(snap, snap)
    assert delta["verdict"] == "CLEAN"
    assert delta["verdict_display"] == "NO DELTA REGRESSION OBSERVED"
    assert delta["verdict_scope"] == "delta_only"
    assert "not a cutover authorization" in delta["verdict_note"].lower()

    out = tmp_path / "delta.xlsx"
    write_diff_workbook(
        snap,
        snap,
        str(out),
        precert={
            "verdict": "CONDITIONAL",
            "verdict_note": "one named blind spot remains",
            "flows": {},
            "stamps": {},
            "segmentation": [],
            "intents": [],
            "blind_spots": ["unverified route surface"],
        },
    )
    from openpyxl import load_workbook

    wb = load_workbook(out, read_only=True)
    try:
        rows = list(wb["Summary"].iter_rows(values_only=True))
    finally:
        wb.close()
    gate = next(row for row in rows if row[0] == "CUTOVER GATE VERDICT")
    observation = next(row for row in rows if row[0] == "DELTA OBSERVATION")
    assert gate[2] == "CONDITIONAL"
    assert observation[2] == "NO DELTA REGRESSION OBSERVED"


def test_certificate_binds_exact_hashes_schema_and_content_semantics():
    cert = precert.compute_precert(
        {"script_version": "V1"},
        {"script_version": "V2"},
        source_hashes={"before": "sha256:" + "a" * 64,
                       "after": {"sha256": "sha256:" + "b" * 64}},
        schema_status={"status": "mismatch", "message": "V1 versus V2", "override": False},
    )

    assert cert["verdict"] == "FAIL"
    assert cert["source_binding"] == {
        "before": "sha256:" + "a" * 64,
        "after": "sha256:" + "b" * 64,
    }
    assert cert["stamps"]["before"]["snapshot_content_sha256"].startswith("sha256:")
    assert cert["schema_status"]["status"] == "mismatch"
    assert any("schema mismatch" in item for item in cert["gate_failures"])


def test_readiness_freeze_with_failed_phase_never_certifies_ready_and_keeps_binding():
    snap = {
        "script_version": "V1",
        "migration_readiness": [{"group": "G1", "readiness": "READY", "checks": []}],
        "assessment_integrity": {"failed_phases": ["Migration Readiness"]},
    }
    cert = precert.compute_readiness_freeze(
        snap,
        source_hash="sha256:" + "c" * 64,
        schema_status=("ok", ""),
    )
    assert cert["verdict"] == "INDETERMINATE"
    assert cert["source_binding"]["snapshot"] == "sha256:" + "c" * 64
    assert cert["integrity"]["ok"] is False


def test_ecmp_partial_blackhole_is_a_hard_delta_and_certificate_gate(monkeypatch):
    partial = {
        "src": "10.0.1.1",
        "dst": "10.0.2.1",
        "ecmp_dropping_legs": [{"next_hop": "10.0.12.2"}],
    }

    def fake_delta(*_args, **_kwargs):
        return {
            "assessed": True,
            "preserved": 1,
            "both_unreachable": 0,
            "newly_blocked": [],
            "newly_reachable": [],
            "inconclusive_pairs": [],
            "ecmp_partial_drop": [partial],
            "pairs_tested": 1,
            "subnets_tested": 2,
            "subnets_total": 2,
            "capped": False,
        }

    monkeypatch.setattr("cisco_toolkit.fib.reachability_delta", fake_delta)
    snap = _comparable_snapshot(segmentation={"domains": [], "vrfs": []})

    delta = compute_snapshot_delta(snap, snap)
    cert = precert.compute_precert(snap, snap)

    assert delta["verdict"] == "REGRESSED"
    assert cert["verdict"] == "FAIL"
    assert any("partial packet loss" in item for item in cert["gate_failures"])


def test_duplicate_and_malformed_path_intents_remain_visible_and_fail_closed():
    intents = [
        {"id": "dup", "src": "10.0.0.1", "dst": "10.0.0.2", "expect": "REACHES"},
        {"id": "dup", "src": "not-an-ip", "dst": "10.0.0.3", "expect": "REACHES"},
    ]
    one = path_assertions.evaluate_path_assertions({}, intents)
    pair = path_assertions.revalidate({}, {}, intents)
    cert = precert.compute_precert({}, {}, path_intents=intents)

    assert one["summary"]["invalid"] == 2
    assert len(one["results"]) == len(pair["results"]) == 2
    assert all(row["verdict"] == "fail" and not row["valid"] for row in one["results"])
    assert pair["summary"]["invalid"] == 2
    assert cert["verdict"] == "FAIL"
    assert len([x for x in cert["gate_failures"] if "path intent" in x]) == 2


def _fhrp_snapshot():
    return {
        "routes": {},
        "stp_roots": {},
        "fhrp_detail": {
            "core1": [{
                "ifname": "Vlan10", "group": "10", "state": "Active", "priority": 110,
                "vip": "10.0.10.1", "version": 2,
            }],
            "dist1": [{
                "ifname": "Vlan10", "group": "10", "state": "Standby", "priority": 100,
                "vip": "10.0.10.1", "version": 2,
            }],
        },
    }


def test_cutover_simulator_rejects_unknown_fhrp_target_without_mutation():
    snap = _fhrp_snapshot()
    original = deepcopy(snap)
    result = cutover_sim.simulate_cutover(
        snap,
        [{"action": "move_fhrp_active", "ifname": "Vlan10", "group": "10",
          "to_host": "not-a-member"}],
        pairs=[],
    )
    row = result["steps"][0]

    assert row["valid"] is False and row["is_noop"] is True
    assert "not a member" in " ".join(row["validation_errors"])
    assert row["fhrp_takeovers"] == []
    assert result["summary"]["n_invalid_steps"] == 1
    assert snap == original


def test_cutover_simulator_rejects_ambiguous_default_fhrp_target():
    snap = _fhrp_snapshot()
    snap["fhrp_detail"]["dist2"] = [{
        "ifname": "Vlan10", "group": "10", "state": "Listen", "priority": 100,
        "vip": "10.0.10.1", "version": 2,
    }]
    result = cutover_sim.simulate_cutover(
        snap,
        [{"action": "move_fhrp_active", "ifname": "Vlan10", "group": "10"}],
        pairs=[],
    )
    assert result["steps"][0]["valid"] is False
    assert "ambiguous" in " ".join(result["steps"][0]["validation_errors"])


def test_mixed_known_and_unknown_rules_never_pass():
    result = assertions.evaluate_assertion(
        {"value": "hello"},
        {
            "id": "mixed",
            "subject": "value",
            "all_of": [
                {"type": "contains", "value": "hello"},
                {"type": "typo-rule", "value": "ignored"},
            ],
        },
    )
    assert result["status"] == assertions.NOT_OBSERVED
    assert "PASS withheld" in result["detail"]


def test_per_object_assertions_are_wired_into_pack_results_and_fail_on_duplicates():
    snap = {
        "interfaces": {
            "sw1": {"Gi1": {"mtu": 1500, "ip": "10.0.0.1"}},
            "sw2": {"Gi1": {"mtu": 1500, "ip": "10.0.0.1"}},
        }
    }
    pack = {
        "for_each": [{
            "id": "interfaces",
            "collection": "interfaces",
            "field_rules": [{"field": "mtu", "op": "max", "value": 1500}],
            "unique_by": "ip",
        }]
    }
    result = assertions.evaluate_pack(snap, pack)

    assert result["summary"]["grade"] == "fail"
    assert len(result["object_results"]) == 1
    object_result = result["object_results"][0]
    assert object_result["status"] == assertions.FAIL
    assert object_result["object_evaluation"]["summary"]["n_uniqueness_violations"] == 1


def test_missing_per_object_collection_is_exposed_as_not_observed():
    result = assertions.evaluate_pack(
        {},
        {"for_each": [{"id": "interfaces", "collection": "interfaces",
                       "field_rules": [{"field": "mtu", "op": "required"}]}]},
    )
    assert result["summary"]["grade"] == "na"
    assert result["object_results"][0]["status"] == assertions.NOT_OBSERVED


def test_constraintless_per_object_declaration_cannot_pass_vacuously():
    result = assertions.evaluate_pack(
        {"interfaces": [{"host": "sw1", "mtu": 1500}]},
        {"for_each": [{"id": "empty", "collection": "interfaces"}]},
    )
    assert result["object_results"][0]["status"] == assertions.NOT_OBSERVED
    assert result["summary"]["grade"] == "na"


def test_fhrp_consistency_requires_forwarding_and_backup_roles_and_scopes_domains():
    from cisco_toolkit.excel import compute_fhrp_consistency
    from cisco_toolkit.model import InterfaceData

    no_active = {
        "a": {"Vlan10": InterfaceData(
            port="Vlan10", svi_ip="10.0.10.2 255.255.255.0",
            hsrp_behavior="HSRP grp 10 Standby VIP 10.0.10.1")},
        "b": {"Vlan10": InterfaceData(
            port="Vlan10", svi_ip="10.0.10.3 255.255.255.0",
            hsrp_behavior="HSRP grp 10 Listen VIP 10.0.10.1")},
    }
    rows = compute_fhrp_consistency(no_active)
    assert any("no observed Active/Master" in issue for issue in rows[0]["issues"])

    separate_domains = {
        "a": {"Vlan10": InterfaceData(
            port="Vlan10", svi_ip="10.0.10.2 255.255.255.0", vrf="RED",
            hsrp_behavior="HSRP grp 10 Active VIP 10.0.10.1")},
        "b": {"Vlan10": InterfaceData(
            port="Vlan10", svi_ip="10.0.20.2 255.255.255.0", vrf="BLUE",
            hsrp_behavior="HSRP grp 10 Active VIP 10.0.20.1")},
    }
    assert compute_fhrp_consistency(separate_domains) == []


def test_excel_surfaces_protocol_and_executive_integrity_failures():
    from openpyxl import Workbook
    from cisco_toolkit.excel import (
        EXEC_SUMMARY_SHEET_NAME,
        PROTOCOL_INTELLIGENCE_SHEET_NAME,
        write_executive_summary_sheet,
        write_protocol_intelligence_sheet,
    )

    wb = Workbook()
    write_protocol_intelligence_sheet(wb, [], unavailable=True)
    assert "UNVERIFIED" in wb[PROTOCOL_INTELLIGENCE_SHEET_NAME]["A2"].value

    write_executive_summary_sheet(
        wb,
        [],
        [],
        [],
        [],
        brief={},
        assessment_integrity={
            "failed_phases": ["Health Scores", "Migration Punch-List", "Migration Readiness"]
        },
    )
    values = [
        str(cell.value)
        for row in wb[EXEC_SUMMARY_SHEET_NAME].iter_rows()
        for cell in row
        if cell.value is not None
    ]
    text = "\n".join(values)
    assert "ASSESSMENT INTEGRITY - UNVERIFIED" in text
    assert "Migration punch-list — UNVERIFIED" in text
    assert "Migration readiness (per move group) — UNVERIFIED" in text
    assert "Critical band\nUNVERIFIED" in text


def _document_text(document) -> str:
    parts = [p.text for p in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            parts.extend(cell.text for cell in row.cells)
    return "\n".join(parts)


def test_docx_consumers_withhold_clean_claims_on_failed_empty_sections(tmp_path):
    pytest.importorskip("docx")
    from docx import Document
    from cisco_toolkit.design import write_design_doc_docx
    from cisco_toolkit.engagement import write_engagement_docx
    from cisco_toolkit.runbook import write_runbook_docx

    snap = {
        "devices": {},
        "interfaces": {},
        "health_scores": [],
        "punchlist": [],
        "migration_readiness": [],
        "lifecycle_risk": {},
        "design_blueprint": {},
        "assessment_integrity": {
            "failed_phases": [
                "Health Scores", "Migration Punch-List", "Migration Readiness",
                "Lifecycle risk", "Design blueprint",
            ]
        },
    }
    outputs = [
        (write_design_doc_docx, tmp_path / "design.docx"),
        (write_runbook_docx, tmp_path / "runbook.docx"),
        (write_engagement_docx, tmp_path / "engagement.docx"),
    ]
    for writer, path in outputs:
        writer(str(path), deepcopy(snap), "Integrity test")
        text = _document_text(Document(path))
        assert "UNVERIFIED" in text
        assert "empty fallback" in text.lower()
    engagement_text = _document_text(Document(tmp_path / "engagement.docx"))
    assert "HOLD" in engagement_text
    design_text = _document_text(Document(tmp_path / "design.docx"))
    assert "No punch-list items" not in design_text


def test_deck_withholds_green_schedule_claim_when_any_phase_failed(tmp_path):
    pytest.importorskip("pptx")
    from pptx import Presentation
    from cisco_toolkit.deck import write_executive_deck_pptx

    out = tmp_path / "deck.pptx"
    write_executive_deck_pptx(
        str(out),
        {
            "health_scores": [],
            "punchlist": [],
            "migration_readiness": [],
            "assessment_integrity": {"failed_phases": ["Migration Punch-List"]},
        },
        "Integrity test",
    )
    prs = Presentation(out)
    text = "\n".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert "ASSESSMENT UNVERIFIED" in text
    assert "fleet is in good shape to schedule" not in text
